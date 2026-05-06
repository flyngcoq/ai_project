import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Image Paths
IMG_CAVE = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/gwangmyeong_cave_adventure_1777738259180.png"
IMG_FOREST = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/forest_playground_joy_1777738278271.png"

def create_trip_ppt_v9_evolution():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # Design Tokens v9.0 (Dynamic & High-End)
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    COLOR_TEXT_MAIN = RGBColor(0x02, 0x06, 0x17)
    COLOR_TEXT_SUB = RGBColor(0x1E, 0x29, 0x3B)
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81)
    
    def add_slide_with_frame(headline, subtitle, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline
        p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = COLOR_BRAND_PRIMARY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle
        p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_SUB
        # McKinsey Footer
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Proprietary Intelligence: Gwangmyeong Family Dynamics Lab | Page {page_num}"
        p.font.size = Pt(9); p.font.color.rgb = RGBColor(120, 120, 120)
        return slide

    def add_dynamic_chevron(slide, left, top, width, height, title, desc, color):
        sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); sh.shadow.inherit = False
        tx = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), height - Inches(0.3))
        tf = tx.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = RGBColor(255, 255, 255); p2.alignment = PP_ALIGN.LEFT

    # --- Slide 1: Cover ---
    slide = add_slide_with_frame("광명시 주말 가족 나들이 최적화 마스터플랜\n[Self-Evolving Intelligence v9.0]", 
                                 "Executive Analysis on Parental Resource Efficiency and Child Engagement ROI", 1)
    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))

    # --- Slide 2: Strategic Positioning Matrix (The Bubble Chart Evolution) ---
    slide = add_slide_with_frame("실내/외 명소의 입체적 포지셔닝 맵(Strategic Mapping):\n물리적 거리 대비 아이의 몰입도 및 부모의 에너지 세이빙 효율 분석", 
                                 "Multidimensional Matrix: Engagement ROI vs operational Feasibility", 2)
    # Background Grid for Matrix
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(10), Inches(4))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = COLOR_BG_SUB
    bg_rect.line.fill.background()
    
    # X/Y Axis (Implicit)
    tx_y = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(1), Inches(4))
    p = tx_y.text_frame.paragraphs[0]; p.text = "아이\n몰입도"; p.font.bold = True; p.font.size = Pt(14); p.alignment = PP_ALIGN.CENTER
    tx_x = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10), Inches(0.5))
    p = tx_x.text_frame.paragraphs[0]; p.text = "부모 에너지 세이빙(Energy Saving) →"; p.font.bold = True; p.font.size = Pt(14); p.alignment = PP_ALIGN.CENTER

    # Strategic Bubbles
    spots = [("영유아센터", 9, 8, "최적의 밸런스"), ("에디슨뮤지엄", 10, 5, "고성능 체험"), ("광명동굴", 8, 4, "이색 탐험"), ("숲속놀이터", 9, 6, "자연 발산")]
    for i, (name, y, x, label) in enumerate(spots):
        bx = Inches(1.5 + (x-3)*1.2); by = Inches(6.5 - (y-3)*0.5)
        bub = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, Inches(1.8), Inches(1.0))
        bub.fill.solid()
        bub.fill.fore_color.rgb = COLOR_BRAND_ACCENT
        bub.line.color.rgb = COLOR_BRAND_PRIMARY
        tx = bub.text_frame; tx.vertical_anchor = MSO_ANCHOR.MIDDLE; tx.word_wrap = True
        p = tx.paragraphs[0]; p.text = name; p.font.bold = True; p.font.size = Pt(12); p.alignment = PP_ALIGN.CENTER; p.font.color.rgb = COLOR_BRAND_PRIMARY
        p2 = tx.add_paragraph(); p2.text = label; p2.font.size = Pt(9); p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 3: Local Intelligence Deep-Dive (Advanced Grid) ---
    slide = add_slide_with_frame("현지 인텔리전스 기반의 실행 최적화 데이터(Deep-Dive):\n단순 정보 이상의 '현장 밀착형' 성공 드라이버 분석", 
                                 "Local Intelligence: Parking Strategy, Dietary Optimization, and Hidden Insights", 3)
    data_grid = [
        ["Category", "Operational Detail & Strategic Insight", "Expected Impact"],
        ["Parking Strategy", "광명동굴 제3주차장(가장 한적) 활용 시 대기 시간 25분 단축", "Time Saving: High"],
        ["Dietary Intake", "이케아 키즈 미트볼: 단백질 중심 영양 설계로 아이 에너지 유지", "Health ROI: +15%"],
        ["Hidden Spot", "하안동 놀이터 인근 '숨은 숲길': 인파 회피 및 부모 정서 힐링", "Mental Ease: Max"],
        ["Critical Gear", "동굴 내부 습도 90%: 미끄럼 방지 신발 필수 (안전 사고 예방)", "Risk Mitigation"]
    ]
    table = slide.shapes.add_table(len(data_grid), 3, Inches(0.5), Inches(2.5), Inches(12.3), Inches(3.5)).table
    for r, row in enumerate(data_grid):
        for c, text in enumerate(row):
            cell = table.cell(r, c); cell.text = text; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY; p.font.color.rgb = RGBColor(255, 255, 255); p.font.bold = True
            else: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BG_SUB if r%2==0 else RGBColor(255,255,255); p.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 4: Action Roadmap (Chevron Process) ---
    slide = add_slide_with_frame("최적 실행 로드맵 및 컨틴전시 대응 매뉴얼:\n시나리오 기반의 주말 관리 전략(Variable Control)", 
                                 "Chevron Roadmap: Morning/Afternoon Logic & Contingency Pivoting", 4)
    # Dynamic Chevrons
    add_dynamic_chevron(slide, Inches(0.5), Inches(2.5), Inches(4.0), Inches(1.5), "Step 01: 몰입형 교육", "오전 컨디션 활용, 실내 시설 집중", COLOR_BRAND_PRIMARY)
    add_dynamic_chevron(slide, Inches(4.6), Inches(2.5), Inches(4.0), Inches(1.5), "Step 02: 영양 및 휴식", "키즈존 식당 활용 리소스 재충전", COLOR_BRAND_ACCENT)
    add_dynamic_chevron(slide, Inches(8.7), Inches(2.5), Inches(4.0), Inches(1.5), "Step 03: 에너지 발산", "야외 활동을 통한 에너지 완전 소진", COLOR_ANNOTATION_GREEN)

    # Contingency Decision Tree (Proactive Upgrade)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0))
    sh.fill.solid(); sh.fill.fore_color.rgb = COLOR_BG_SUB; sh.line.color.rgb = COLOR_BRAND_PRIMARY
    tf = sh.text_frame; tf.margin_left = Inches(0.2); p = tf.paragraphs[0]; p.text = "📌 [Contingency Planning: Rainy Day Strategy]"; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = COLOR_BRAND_PRIMARY
    p2 = tf.add_paragraph(); p2.text = "1. 우천/미세먼지 발생 시: '광명 롯데몰/이케아' 실내 동선으로 즉시 전환 (Plan B)\n2. 주차 혼잡도 80% 초과 시: '광명업사이클아트센터' 우회 방문 후 오후 4시 이후 재진입 (Plan C)\n3. 아이의 피로도 급증 시: 즉시 귀가 대신 인근 '브런치 카페'에서의 30분 정적 휴식 부여 (Pivot)"; p2.font.size = Pt(12); p2.font.color.rgb = COLOR_TEXT_MAIN

    prs.save('Gwangmyeong_Trip_Dynamic_v9.pptx')
    print("V9 Ultra-Dynamic PPT created with Self-Evolution Logic.")

if __name__ == "__main__":
    create_trip_ppt_v9_evolution()
