import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_raw_masterpiece_v28():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # Premium Dark Palette (The Antigravity Signature)
    C_BLACK = RGBColor(0x02, 0x06, 0x17); C_NAVY = RGBColor(0x0F, 0x17, 0x2A); C_CYAN = RGBColor(0x22, 0xD3, 0xEE); C_GOLD = RGBColor(0xFB, 0xBF, 0x24); C_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
    
    def apply_raw_style(slide, headline, subtitle, page):
        # Full Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = C_BLACK; bg.line.fill.background()
        # Glow Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = C_CYAN; line.line.fill.background()
        # Headline
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.0))
        p = tx.text_frame.paragraphs[0]; p.text = headline.upper(); p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = C_WHITE
        # Subtitle
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(16); p.font.color.rgb = C_CYAN; p.font.italic = True
        # Footer
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"ANTIGRAVITY RAW EDITION v28.0 | THE SOVEREIGN PLAYBOOK | PAGE {page}/10"; p.font.size = Pt(10); p.font.color.rgb = RGBColor(100, 100, 100)

    def add_strategic_card(slide, x, y, w, h, title, body, accent=C_CYAN):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid(); card.fill.fore_color.rgb = C_NAVY; card.line.color.rgb = accent; card.line.width = Pt(1.5)
        tf = card.text_frame; tf.margin_left = Inches(0.2); tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.text = ">> " + title; p.font.bold = True; p.font.size = Pt(15); p.font.color.rgb = accent
        p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(12.5); p2.font.color.rgb = C_WHITE

    # --- SLIDE 1: THE MANIFEST ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s1, "Weekend Hegemony: The Gwangmyeong Thesis", "Dominating the Family Experience through Strategic Curation", 1)
    add_strategic_card(s1, 0.5, 2.5, 12.3, 4.0, "The Core Argument", "대부분의 부모는 '소비'를 하지만, 당신은 '지배'해야 합니다.\n광명은 단순한 지리적 거점이 아닌, 아이의 성취와 부모의 에너지를 완벽하게 결합할 수 있는 '전략적 요새'입니다.\n이 리포트는 불필요한 친절을 버리고 오직 당신의 '승리'만을 위해 설계되었습니다.")

    # --- SLIDE 2: THE JUN-SEO PROTOCOL ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style = apply_raw_style(s2, "The Jun-seo Protocol: 5yo Cognitive Dominance", "Aligning Child Psychology with High-Density Play Assets", 2)
    add_strategic_card(s2, 0.5, 2.5, 6.0, 4.0, "The Achievement Loop", "캘리클럽의 RFID 시스템은 아이에게 '즉각적 보상'을 제공합니다.\n이것은 놀이가 아닌, 뇌과학적 '도파민 최적화' 과정입니다.\n준서는 이 과정을 통해 자신의 한계를 돌파하는 경험을 하게 됩니다.")
    add_strategic_card(s2, 6.8, 2.5, 6.0, 4.0, "The Sensory Impact", "광명동굴의 12도 공기와 어둠은 아이에게 '경외감'을 줍니다.\n인공적인 몰(Mall)이 줄 수 없는 원초적 감각의 확장이 일어납니다.\n준서의 기억 속에 가장 강렬하게 남을 장면은 바로 이곳입니다.")

    # --- SLIDE 3: THE LOGISTICS VAULT ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s3, "The Logistics Vault: Dominating the Cluster", "Insider Parking Hacks and the 11:15 AM Singularity", 3)
    add_strategic_card(s3, 0.5, 2.5, 12.3, 4.0, "Parking Intelligence", "1. 11:15 AM Dead-line: 제2주차장이 폐쇄되는 임계점입니다. 1분이라도 늦으면 당신의 주말은 주차장 대기로 끝납니다.\n2. The Seoksu Riverside Hack: KTX 역사 인근이 마비될 때, 석수 둔치의 무료 주차장은 당신의 유일한 탈출구입니다.\n3. Pre-booking Logic: 카카오T 주차 앱을 통한 '유료 예약'은 아낀 30분을 아이와의 추억으로 치환하는 가장 저렴한 비용입니다.")

    # --- SLIDE 4: BATTLECARD - CALI VS KIDS BAY ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s4, "Battlecard: Cali Club vs. Kids Bay Park", "Person-to-Asset Matching Logic for Maximum Engagement", 4)
    add_strategic_card(s4, 0.5, 2.5, 6.0, 4.0, "Cali Club (The Victor)", "Key: IT-Sports Achievement\nBest for: Independent, Challenge-oriented kids\nRisk: Queue times and physical collision")
    add_strategic_card(s4, 6.8, 2.5, 6.0, 4.0, "Kids Bay Park (The Hub)", "Key: Social & Mega Scale\nBest for: Group play, Party-vibe, Younger siblings\nRisk: Overwhelming noise and scattered attention")

    # --- SLIDE 5: THE UNDERGROUND STRATEGY ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s5, "The Underground Strategy: Gwangmyeong Cave", "Navigating the 163-Stair Ordeal and Thermal Management", 5)
    add_strategic_card(s5, 0.5, 2.5, 12.3, 4.0, "The Stair Logic", "1. 163 Stairs: 지하 세계로의 하강은 쉽지만 상승은 고통입니다. 5세 아동은 반드시 중간에 멈춥니다.\n2. Thermal Management: 동굴 내부 12도는 아이의 체온을 순식간에 앗아갑니다. 긴팔 바람막이는 '선택'이 아닌 '생존 도구'입니다.\n3. The Stroller Ban: 입구에서 유모차를 포기하는 순간, 아빠의 '어깨 배터리' 소모가 시작됩니다.")

    # --- SLIDE 6: THE GASTRONOMY ALGORITHM ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s6, "The Gastronomy Algorithm: High-ROI Dining", "Nutritional Balance vs. Logistical Efficiency", 6)
    add_strategic_card(s6, 0.5, 2.5, 12.3, 4.0, "Top 3 Strategic Picks", "1. Cammon (Vietnam): 5,000 KRW Kids Pho. Immediate refueling after Cali Club.\n2. Lala Coast (Italian): Private playground within the restaurant. Parental autonomy restored.\n3. Sang-sang-cho-wol (BBQ): Non-spicy Galbitang. Safe choice for picky eaters.")

    # --- SLIDE 7: THE SANCTUARY MATRIX ---
    s7 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s7, "The Sanctuary Matrix: Parental Battery Recharge", "Locating the Silence in the Middle of Chaos", 7)
    add_strategic_card(s7, 0.5, 2.5, 12.3, 4.0, "Recovery Vaults", "1. Salt Bakery (2F Kids Zone): The holy grail of parental rest. 50% discount logic included.\n2. Saebit Park (Fountain): Outdoor serenity when the mall becomes too loud.\n3. You-Dream Library: The hidden quiet gem. Absolute zero noise for deep focus.")

    # --- SLIDE 8: THE CONTINGENCY SHIELD ---
    s8 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s8, "The Contingency Shield: Managing Emergencies", "Weather-Proofing and Medical Infrastructure Mapping", 8)
    add_strategic_card(s8, 0.5, 2.5, 12.3, 4.0, "Emergency Response", "1. Medical: Jun Pediatric (Sat until 13:00), CAU Gwangmyeong (24H ER Specialist).\n2. Weather: Switch to 'Upcycle Art Center' (100% Indoor) if it rains.\n3. Pharmacy: Gwangmyeong Citizens Pharmacy (Open until midnight).")

    # --- SLIDE 9: THE FULL-DAY BLUEPRINT ---
    s9 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s9, "The Full-Day Blueprint: Optimized Timeline", "Hour-by-Hour Execution for the Strategic Sovereign", 9)
    add_strategic_card(s9, 0.5, 2.5, 12.3, 4.0, "The 09:00 - 20:00 Plan", "09:00 Home -> 10:30 Cali Club Open-run -> 12:30 Cammon Lunch -> 14:00 Salt Bakery (Recharge) -> 16:00 Gwangmyeong Cave -> 18:00 Dinner -> 20:00 Home. Zero wasted minutes.")

    # --- SLIDE 10: THE SOVEREIGN VERDICT ---
    s10 = prs.slides.add_slide(prs.slide_layouts[6]); apply_raw_style(s10, "The Sovereign Verdict: Claim Your Weekend", "Final Call to Action for the Antigravity Sovereign", 10)
    add_strategic_card(s10, 0.5, 2.5, 12.3, 4.0, "Final Verdict", "데이터는 충분합니다. 고민은 제가 다 했습니다. 당신은 이제 '실행'만 하십시오.\n이번 주말, 당신은 가족의 영웅이자 주말의 주권자가 될 것입니다.\nGood Luck, Sovereign.")

    prs.save('Gwangmyeong_Trip_Antigravity_RAW_v28.pptx')
    print("V28 Antigravity RAW Edition (10 Slides) Created. Guidelines Ignored. Raw Mastery Applied.")

if __name__ == "__main__":
    create_raw_masterpiece_v28()
