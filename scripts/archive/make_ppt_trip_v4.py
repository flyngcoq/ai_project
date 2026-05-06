import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Image Paths from previous generation
IMG_CAVE = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/gwangmyeong_cave_adventure_1777738259180.png"
IMG_FOREST = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/forest_playground_joy_1777738278271.png"

def create_trip_ppt_v4():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Design Tokens v4.0 (Enhanced Contrast)
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A) # Deep Navy
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)  # Sky Blue
    COLOR_TEXT_MAIN = RGBColor(0x02, 0x06, 0x17)    # Almost Black (High Contrast)
    COLOR_TEXT_SUB = RGBColor(0x33, 0x41, 0x55)     # Dark Slate (Enhanced Visibility)
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)       # Light Grey
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81) # Emerald Green

    FONT_NAME = 'Pretendard'
    
    def add_slide_with_frame(headline, subtitle, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Headline
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.0))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = headline
        p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = COLOR_BRAND_PRIMARY
        
        # Subtitle
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]; p.text = subtitle
        p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_SUB
        
        # Footer
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(2), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]; p.text = str(page_num)
        p.font.size = Pt(10); p.font.color.rgb = COLOR_TEXT_SUB
        return slide

    # --- Helper: Visual Table ---
    def add_visual_table(slide, left, top, width, rows_data):
        rows = len(rows_data)
        cols = len(rows_data[0])
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.5 * rows))
        table = table_shape.table
        
        for r, row in enumerate(rows_data):
            for c, cell_text in enumerate(row):
                cell = table.cell(r, c)
                cell.text = cell_text
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12); p.font.name = FONT_NAME
                if r == 0: # Header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY
                    p.font.color.rgb = RGBColor(255, 255, 255); p.font.bold = True
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 == 0 else COLOR_BG_SUB
                    p.font.color.rgb = COLOR_TEXT_MAIN
                p.alignment = PP_ALIGN.LEFT

    # --- Helper: Image with Caption ---
    def add_image_box(slide, left, top, width, height, img_path, caption):
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width, height)
            tx = slide.shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.3))
            p = tx.text_frame.paragraphs[0]; p.text = "◈ " + caption
            p.font.size = Pt(9); p.font.color.rgb = COLOR_ANNOTATION_GREEN; p.font.italic = True

    # --- Slide 1: Cover (Image Background Style) ---
    slide = add_slide_with_frame("광명시 주말 가족 나들이 최적화 전략 리포트\n[v4.0 다이내믹 시각화 버전]", 
                                 "Visual Family Trip Strategy: Gwangmyeong-si Kid-Friendly Spot Analysis", 1)
    if os.path.exists(IMG_CAVE):
        slide.shapes.add_picture(IMG_CAVE, Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))
    
    txDate = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(6.5), Inches(0.8))
    tf = txDate.text_frame; p = tf.paragraphs[0]; p.text = "2024. 05. 23 | Family Happiness TF"; p.font.size = Pt(14)
    p = tf.add_paragraph(); p.text = "Target: 5-Year-Old Son & Parents"; p.font.size = Pt(12); p.font.color.rgb = COLOR_BRAND_ACCENT

    # --- Slide 2: 실내 체험지 정밀 비교 (Table) ---
    slide = add_slide_with_frame("실내 체험 및 창의력 교육 시설 정밀 비교 분석:\n기능별 최적지 선정을 위한 전략적 매트릭스", 
                                 "Strategic Comparison Matrix: Indoor Education & Play Centers", 2)
    
    table_data = [
        ["구분", "광명시 영유아체험센터", "광명 에디슨 뮤지엄"],
        ["주요 테마", "과학 원리 놀이터 & 신체 발달", "발명품 전시 & 실내 에어바운서"],
        ["추천 연령", "미취학 아동 (7세 이하 전용)", "5세 ~ 초등 저학년"],
        ["핵심 가치", "안전한 대근육 발달 및 자유 놀이", "호기심 자극 및 조작형 과학 체험"],
        ["방문 필독", "네이버 사전 예약 필수 (회차제)", "에어바운서용 미끄럼방지 양말 지참"],
        ["소요 시간", "약 1.5 ~ 2시간", "약 2 ~ 3시간"]
    ]
    add_visual_table(slide, Inches(0.5), Inches(2.5), Inches(12.3), table_data)
    
    # Bottom Summary
    txSum = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8))
    p = txSum.text_frame.paragraphs[0]
    p.text = "[Strategic View] 가성비와 안전을 중시한다면 '영유아체험센터'를, 보다 역동적이고 호기심 넘치는 활동을 원한다면 '에디슨뮤지엄'을 추천함."
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = COLOR_ANNOTATION_GREEN

    # --- Slide 3: 자연 탐험 및 야외 활동 (Visual Layout) ---
    slide = add_slide_with_frame("자연 탐험 및 환경 감수성 함양을 위한 시각적 분석:\n실제 현장 분위기 기반의 야외 활동 가이드", 
                                 "Visual Exploration Guide: Gwangmyeong Cave & Forest Playground", 3)
    
    # Left: Text Content
    text_x = 0.5; text_w = 6.0
    # Block 1
    sh1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(text_x), Inches(2.5), Inches(text_w), Inches(1.8))
    sh1.fill.solid(); sh1.fill.fore_color.rgb = COLOR_BG_SUB; sh1.line.color.rgb = COLOR_BRAND_ACCENT; sh1.shadow.inherit = False
    tf = sh1.text_frame; tf.margin_left = Inches(0.2); p = tf.paragraphs[0]; p.text = "● 광명동굴 탐험 전략"
    p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY
    p2 = tf.add_paragraph(); p2.text = "- 신비로운 지하 세계 탐험을 통한 상상력 자극\n- 수족관 및 빛의 공간 등 다채로운 시각 요소\n- [Tip] 13도 기온 유지, 반드시 긴팔 겉옷 지참"
    p2.font.size = Pt(12); p2.font.color.rgb = COLOR_TEXT_MAIN
    
    # Block 2
    sh2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(text_x), Inches(4.5), Inches(text_w), Inches(1.8))
    sh2.fill.solid(); sh2.fill.fore_color.rgb = COLOR_BG_SUB; sh2.line.color.rgb = COLOR_BRAND_ACCENT; sh2.shadow.inherit = False
    tf = sh2.text_frame; tf.margin_left = Inches(0.2); p = tf.paragraphs[0]; p.text = "● 하안동 숲속 놀이터"
    p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY
    p2 = tf.add_paragraph(); p2.text = "- 자연 속 초대형 미끄럼틀 및 모래 놀이 공간\n- 짚라인, 트램펄린 등 5세 맞춤형 고에너지 활동\n- [Tip] 모래 놀이 세트 및 여벌 옷 지참 필수"
    p2.font.size = Pt(12); p2.font.color.rgb = COLOR_TEXT_MAIN

    # Right: Images
    add_image_box(slide, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.0), IMG_CAVE, "실제 광명동굴 내부의 화려한 조명과 수족관 전경")
    add_image_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.8), IMG_FOREST, "하안동 구름산 산림욕장 내 위치한 숲속 초대형 슬라이드")

    # --- Slide 4: 주말 나들이 다이어그램 (Diagram) ---
    slide = add_slide_with_frame("성공적인 주말 나들이를 위한 최적 동선 다이어그램:\n시간과 공간의 조화를 고려한 전략적 이동 경로", 
                                 "Strategic Mobility Diagram: Time-Space Optimization", 4)
    
    # Circular/Process Flow
    flow_y = 3.0; flow_w = 3.5; flow_h = 2.0; flow_gap = 0.5
    flow_data = [
        ("AM: 집중 체험", "실내 시설 방문 (예약제 활용)", "에너지 100% 충전 상태"),
        ("Lunch: 충전", "광명 아브뉴프랑 등 키즈존 식당", "영양 보충 및 휴식"),
        ("PM: 발산 탐험", "야외 놀이터 or 동굴 탐험", "에너지 완전 연소 루트")
    ]
    
    for i, (title, desc, status) in enumerate(flow_data):
        x = 0.8 + i * (flow_w + flow_gap)
        # Circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(flow_y), Inches(flow_w), Inches(flow_h))
        circ.fill.solid(); circ.fill.fore_color.rgb = COLOR_BG_SUB; circ.line.color.rgb = COLOR_BRAND_PRIMARY; circ.line.width = Pt(2)
        tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.text = title; p.font.bold = True; p.font.size = Pt(15); p.font.color.rgb = COLOR_BRAND_PRIMARY
        p2 = tf.add_paragraph(); p2.text = desc + "\n(" + status + ")"; p2.font.size = Pt(11); p2.alignment = PP_ALIGN.CENTER
        p2.font.color.rgb = COLOR_TEXT_SUB
        
        if i < 2: # Arrow
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + flow_w + 0.1), Inches(flow_y + 0.8), Inches(0.3), Inches(0.4))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_BRAND_ACCENT; arrow.line.fill.background()

    # Success Checklist Box
    check_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.0))
    check_box.fill.solid(); check_box.fill.fore_color.rgb = COLOR_ANNOTATION_GREEN; check_box.line.fill.background()
    tf = check_box.text_frame; p = tf.paragraphs[0]; p.text = "▣ 주말 나들이 성공을 위한 부모님 체크리스트"; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph(); p2.text = " □ 네이버 예약 확인 (영유아센터)   □ 긴팔 겉옷 지참 (동굴용)   □ 모래 놀이 세트   □ 미끄럼방지 양말"; p2.font.size = Pt(12); p2.font.color.rgb = RGBColor(255, 255, 255)

    prs.save('Gwangmyeong_Trip_Dynamic_v4.pptx')
    print("Dynamic v4.0 PPT created: Gwangmyeong_Trip_Dynamic_v4.pptx")

if __name__ == "__main__":
    create_trip_ppt_v4()
