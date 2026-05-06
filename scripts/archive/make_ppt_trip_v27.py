import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

IMG_CALI = "/Users/flyngcoq/AI_Project/scripts/cali_club_it_sports_vibe_1777742544504.png"
IMG_EDISON = "/Users/flyngcoq/AI_Project/scripts/edison_museum_discovery_vibe_1777742558123.png"

def create_narrative_ppt_v27():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    C_NAVY = RGBColor(0x0F, 0x17, 0x2A); C_BLUE = RGBColor(0x38, 0xBD, 0xF8); C_BG = RGBColor(0xF8, 0xFA, 0xFC); C_ACCENT = RGBColor(0xEF, 0x44, 0x44)
    
    def apply_style(slide, headline, subtitle, page):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.05))
        line.fill.solid(); line.fill.fore_color.rgb = C_BLUE; line.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_NAVY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(14); p.font.color.rgb = RGBColor(100, 100, 100)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Narrative Strategy v27.0 | Jun-seo's Sovereign Saturday | Page {page}/10"; p.font.size = Pt(10); p.font.color.rgb = RGBColor(180, 180, 180)

    def add_narrative_box(slide, x, y, w, h, emotion, risk, verdict):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = C_BG; box.line.color.rgb = C_BLUE; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.margin_left = Inches(0.2)
        p1 = tf.paragraphs[0]; p1.text = "😊 기대 감정: " + emotion; p1.font.size = Pt(12); p1.font.bold = True; p1.font.color.rgb = RGBColor(0, 150, 0)
        p2 = tf.add_paragraph(); p2.text = "⚠️ 현실 리스크: " + risk; p2.font.size = Pt(12); p2.font.bold = True; p2.font.color.rgb = C_ACCENT
        p3 = tf.add_paragraph(); p3.text = "💡 최종 권고: " + verdict; p3.font.size = Pt(13); p3.font.bold = True; p3.font.color.rgb = C_NAVY

    # --- SLIDE 1: COVER ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s1, "광명 주말 내러티브 전략: 5세 준서의 하루", "Mapping Emotions & Risks through Real-World Family Intelligence", 1)

    # --- SLIDE 2: EMOTIONAL ROADMAP ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s2, "준서의 감정 로드맵: 환희와 탄식의 교차점", "The Journey from Morning Excitement to Evening Burnout", 2)
    # Placeholder for a diagram/table
    table = s2.shapes.add_table(5, 3, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.0)).table
    rows = [["Phase", "Key Experience", "Emotional State"], ["Morning", "Cali Club (Achievement)", "High (Excitement)"], ["Noon", "Cammon (Refuel)", "Stable (Satisfaction)"], ["Afternoon", "Gwangmyeong Cave (Exploration)", "Mix (Awe + Fatigue)"], ["Evening", "Salt Bakery (Recharge)", "Neutral (Recovery)"]]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c); cell.text = val; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(13)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = C_NAVY; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True

    # --- SLIDE 3: PHASE 1 (VICTORY ENGINE) ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s3, "Phase 1: '성취의 엔진' 캘리클럽과 대기 지옥", "Achieving Victory while Mitigating the Queue Risk", 3)
    if os.path.exists(IMG_CALI): s3.shapes.add_picture(IMG_CALI, Inches(0.5), Inches(2.2), Inches(6), Inches(4.5))
    add_narrative_box(s3, 6.8, 2.2, 6.0, 4.5, "태그액션 성공 시의 극강의 자존감 향상", "주말 11시 이후 대기 40팀 돌파 (패닉 가능성)", "10:30 오픈런 필수, 실패 시 아브뉴프랑 대체 동선 가동")

    # --- SLIDE 4: PHASE 2 (TACTICAL REFUEL) ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s4, "Phase 2: '미식의 위로' 까몬의 5,000원 어린이 쌀국수", "Synergistic Dining: Refueling the Child's Battery", 4)
    add_narrative_box(s4, 0.5, 2.2, 12.3, 4.2, "캘리클럽 소진 에너지를 고단백 육수로 완벽 보충", "점심 피크 타임(12:30) 웨이팅 및 소음", "캘리클럽 입장 시 미리 캐치테이블 예약 확인")

    # --- SLIDE 5: PHASE 3 (MAJESTIC EXPLORER) ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s5, "Phase 3: '경외감의 탐험' 광명동굴과 163계단", "Visual Awe vs. Physical Exhaustion (The Stroller Crisis)", 5)
    add_narrative_box(s5, 0.5, 2.2, 12.3, 4.2, "지하 세계의 신비로움과 미디어파사드의 전율", "유모차 반입 금지 및 163계단의 체력 부하", "아기띠 지참 필수, 5세 아동은 계단 왕복 시 안아줄 준비 요망")

    # --- SLIDE 6: PHASE 4 (PARENTAL SANCTUARY) ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s6, "Phase 4: '침묵의 충전' 소올투베이커리 명당 사수", "Parental Battery Recovery through Strategic Seating", 6)
    add_narrative_box(s6, 0.5, 2.2, 12.3, 4.2, "창가 숲 뷰와 아이의 기차 놀이가 주는 찰나의 평화", "할인 순서 오류(카페 선방문 시 동굴 50% 할인) 아쉬움", "반드시 카페 선방문하여 동굴 50% 할인권 수령")

    # --- SLIDE 7-10: Additional Logic ---
    s7 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s7, "Comparison: 준서의 눈으로 본 캘리 vs 키즈베이", "Persona-based Decision: IT Victory vs. Mega Play", 7)
    add_narrative_box(s7, 0.5, 2.2, 12.3, 4.2, "내가 영웅이 된 듯한 기분 (성취 지능)", "에너지의 무한 발산과 군중 속의 재미 (사회성)", "준서가 평소 '내가 할게!'라고 자주 한다면 캘리클럽 추천")

    s8 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s8, "The Logistics Risk Matrix", "Parking Saturation & Emergency Medical Support", 8)
    add_narrative_box(s8, 0.5, 2.2, 12.3, 4.2, "11:15 AM: 제2주차장의 최후의 한 칸이 사라지는 시각", "갑작스러운 발열 시: 철산 준소아과(달빛어린이병원) 상시 대기", "주말엔 '모두의주차장' 앱으로 광명역자이타워 보조 주차장 확보")

    s9 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s9, "The Sovereign Decision Algorithm", "Real-time Fatigue-based Path Selection", 9)
    add_narrative_box(s9, 0.5, 2.2, 12.3, 4.2, "부모 배터리 > 80%: 동굴 풀코스 도전", "부모 배터리 < 30%: 소올투베이커리 벙커 행", "내일의 출근 컨디션을 고려하여 17시 이전 귀가 권고")

    s10 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s10, "Final Verdict: 준서의 미소를 위한 최종 권고", "The No-Regret Execution Roadmap", 10)
    add_insight_box = s10.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.2))
    p = add_insight_box.text_frame.paragraphs[0]; p.text = "최종 작전명: [Compacted Victory Path]\n\n1. 10:30: 롯데몰 캘리클럽 오픈런 (성취감 폭발)\n2. 12:30: 아브뉴프랑 까몬 (단백질 충전)\n3. 14:30: 소올투베이커리 2층 명당 (부모 휴식)\n4. 결론: 무리한 동굴 투어보다 '몰 내 콤팩트 동선'이 이번 주말의 주권을 지키는 길입니다."; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = C_NAVY

    prs.save('Gwangmyeong_Trip_Narrative_v27.pptx')
    print("V27 Narrative Strategy PPT (10 Slides) Created. Based on Jun-seo Storyline.")

if __name__ == "__main__":
    create_narrative_ppt_v27()
