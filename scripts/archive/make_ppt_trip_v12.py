import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Image Paths
IMG_CAVE = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/gwangmyeong_cave_adventure_1777738259180.png"
IMG_FOREST = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/forest_playground_joy_1777738278271.png"

def create_trip_ppt_v12_directive():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # Executive Design Tokens
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    COLOR_TEXT_MAIN = RGBColor(0x02, 0x06, 0x17)
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81)
    
    def add_slide_with_frame(headline, subtitle, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = COLOR_BRAND_PRIMARY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_MAIN
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Strategic Execution: Team Leader Approved v12.0 | Page {page_num}"; p.font.size = Pt(9); p.font.color.rgb = RGBColor(120, 120, 120)
        return slide

    def add_directed_shape(slide, shape_type, left, top, width, height, title, body_lines, color=COLOR_BG_SUB):
        sh = slide.shapes.add_shape(shape_type, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.color.rgb = COLOR_BRAND_PRIMARY
        tf = sh.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.margin_left = Inches(0.2); tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.text = "■ " + title; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY
        for line in body_lines:
            p2 = tf.add_paragraph(); p2.text = " - " + line; p2.font.size = Pt(12); p2.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 1: Cover ---
    slide = add_slide_with_frame("광명시 주말 가족 가치 최적화 전략 리포트\n[Team Leader Directive v12.0]", 
                                 "Strategic Resource Allocation & Kid-Parent ROI Optimization through McKinsey 7S", 1)
    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))

    # --- Slide 2: McKinsey 7S Transformation (Strategy Diagnostic) ---
    slide = add_slide_with_frame("맥킨지 7S 기반의 나들이 환경 진단(Diagnostic Matrix):\n경험의 무형 가치(Soft)와 시설의 유형 인프라(Hard) 간의 시너지 분석", 
                                 "7S Framework Analysis: Aligning Experiences, Infrastructure, and Parental Resources", 2)
    # Hard Elements
    add_directed_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(4.0), Inches(3.5), 
                       "Hard Elements (인프라)", ["Strategy: 아이의 몰입형 교육과 부모의 휴식 공존", "Structure: 영유아센터-동굴-베이커리로 이어지는 삼각 클러스터", "Systems: 100% 사전 예약제(영유아센터) 기반의 안정적 동선"])
    # Soft Elements
    add_directed_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.6), Inches(2.5), Inches(4.0), Inches(3.5), 
                       "Soft Elements (역량/가치)", ["Shared Values: 주말 가족 행복 자본의 극대화", "Staff: 전문 에듀케이터가 상주하는 박물관/키즈카페", "Skills: 5세 아이의 대근육 및 정서 발달 자극 역량"], COLOR_BRAND_ACCENT)
    # Shared Value Focus
    add_directed_shape(slide, MSO_SHAPE.OVAL, Inches(8.8), Inches(3.2), Inches(4.0), Inches(2.5), 
                       "Core Value (핵심 가치)", ["Parental Resource Protection", "High-Engagement Child Play", "Operational Efficiency"], COLOR_ANNOTATION_GREEN)

    # --- Slide 3: 2024 Price & Value Comparison (The Real Data) ---
    slide = add_slide_with_frame("실행 가용 자원 및 비용 분석 (2024 Price & Value):\n최신 시장 데이터 기반의 항목별 경쟁 우위 분석", 
                                 "2024 Market Intelligence: Pricing, Amenities, and Competitive Benchmarking", 3)
    table_data = [
        ["항목", "브레드이발소 키즈카페", "에디슨 뮤지엄", "소올투베이커리"],
        ["입장료(아이)", "25,000원 (2시간)", "20,000원 (종일)", "카페 이용 시 키즈룸 무료"],
        ["특이사항", "강력한 캐릭터 IP 몰입", "에어바운서 & 과학체험", "소금빵 및 프리미엄 휴식"],
        ["부모 편의도", "보통 (능동 케어 필요)", "높음 (카페 구역 인접)", "최상 (전용 키즈룸 운영)"],
        ["Expected ROI", "Extreme Happiness", "High Duration Value", "Maximum Parental Rest"]
    ]
    table = slide.shapes.add_table(len(table_data), 4, Inches(0.5), Inches(2.5), Inches(12.3), Inches(3.5)).table
    for r, row in enumerate(table_data):
        for c, t in enumerate(row):
            cell = table.cell(r, c); cell.text = t; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True
            else: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BG_SUB if r%2==0 else RGBColor(255,255,255); p.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 4: Design Thinking Journey (Advanced Visuals) ---
    slide = add_slide_with_frame("디자인 씽킹 기반의 감정 여정 및 솔루션 도출:\nEmpathize에서 Test까지 이어지는 주말 문제 해결 프로세스", 
                                 "User-Centric Design Thinking Journey: Solving Parental Fatigue & Child Boredom", 4)
    # Arrow Flow (Design Thinking Steps)
    dt_steps = [("Empathize (공감)", "부모의 만성 주말 피로 감지"), ("Define (정의)", "안전과 재미의 동시 충족 필요"), ("Ideate (아이디어)", "광명 맞춤형 클러스터 설계"), ("Prototype (실행)", "브레드이발소-소올투 루트"), ("Test (검증)", "아이 행복 지수 및 피로도 측정")]
    for i, (t, d) in enumerate(dt_steps):
        x = 0.5 + i * 2.5
        sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.5), Inches(2.3), Inches(3.0))
        sh.fill.solid(); sh.fill.fore_color.rgb = COLOR_BG_SUB; sh.line.color.rgb = COLOR_BRAND_ACCENT
        tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.4)
        p = tf.paragraphs[0]; p.text = t; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = COLOR_BRAND_PRIMARY
        p2 = tf.add_paragraph(); p2.text = d; p2.font.size = Pt(10); p2.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 5: Strategic Roadmap & Contingency ---
    slide = add_slide_with_frame("최적 실행 타임라인 및 컨틴전시 대응 매뉴얼:\n시나리오 기반의 성공적인 주말 완수를 위한 로드맵", 
                                 "Tactical Roadmap: Time-Based Execution & Variable Control Matrix", 5)
    # Chevron Process
    add_directed_shape(slide, MSO_SHAPE.CHEVRON, Inches(0.5), Inches(2.2), Inches(4.0), Inches(1.8), "AM: 지능 발달", ["에디슨뮤지엄/영유아센터", "오전 집중력 활용", "사전예약 필수"], COLOR_BRAND_PRIMARY)
    add_directed_shape(slide, MSO_SHAPE.CHEVRON, Inches(4.6), Inches(2.2), Inches(4.0), Inches(1.8), "Lunch: 에너지 회복", ["라라코스트/전통시장", "영양 보충 및 부모 휴식", "키즈존 인접 좌석 확보"], COLOR_BRAND_ACCENT)
    add_directed_shape(slide, MSO_SHAPE.CHEVRON, Inches(8.7), Inches(2.2), Inches(4.0), Inches(1.8), "PM: 신체 활동", ["광명동굴/구름산", "잔여 에너지 완전 연소", "동굴 겉옷/놀이터 세면 지참"], COLOR_ANNOTATION_GREEN)

    # Contingency Table (Min 12pt)
    table_c = slide.shapes.add_table(3, 2, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0)).table
    c_data = [["Variable Risk", "Team Leader's Contingency Solution"], ["미세먼지/우천 시", "브레드이발소 키즈카페 실내 올데이 코스로 즉시 전환"], ["아이 컨디션 난조 시", "소올투베이커리 키즈룸 3시간 집중 휴식 및 조기 귀가"]]
    for r, row in enumerate(c_data):
        for c, t in enumerate(row):
            cell = table_c.cell(r, c); cell.text = t; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY; p.font.color.rgb = RGBColor(255,255,255)
            else: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255,255,255); p.font.color.rgb = COLOR_TEXT_MAIN

    prs.save('Gwangmyeong_Trip_Executive_v12.pptx')
    print("V12 Executive Masterpiece Created under Team Leader Directive.")

if __name__ == "__main__":
    create_trip_ppt_v12_directive()
