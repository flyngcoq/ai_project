import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_sovereign_5_ppt_v26():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    C_NAVY = RGBColor(0x0F, 0x17, 0x2A); C_BLUE = RGBColor(0x38, 0xBD, 0xF8); C_BG = RGBColor(0xF8, 0xFA, 0xFC); C_ACCENT = RGBColor(0xEF, 0x44, 0x44)
    
    def apply_style(slide, headline, subtitle, page):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.05))
        line.fill.solid(); line.fill.fore_color.rgb = C_BLUE; line.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = C_NAVY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(15); p.font.color.rgb = RGBColor(100, 100, 100)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Sovereign Verdict v26.0 | Strategic Essence | Page {page}/5"; p.font.size = Pt(10); p.font.color.rgb = RGBColor(180, 180, 180)

    def add_insight_box(slide, x, y, w, h, title, body, color=C_BLUE):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = C_BG; box.line.color.rgb = color; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]; p.text = "🤵 The Core Strategy: " + title; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(13); p2.font.color.rgb = C_NAVY; p2.font.bold = True

    # --- SLIDE 1: THE PSYCHOLOGY OF VICTORY ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s1, "승리의 심리학: 왜 캘리클럽인가?", "Achieving High-Efficacy in 5-Year-Olds through RFID Tag Action", 1)
    add_insight_box(s1, 0.5, 2.2, 12.3, 4.2, "Cognitive ROI Analysis", "1. 성취의 루프: 태그 인지 -> 즉각적 시각 보상 -> 재도전으로 이어지는 긍정적 강화 학습.\n2. 자아 효능감: 단순한 놀이가 아닌 '미션 수행'을 통해 아이는 자신의 유능함을 발견함.\n3. 능동적 주도성: 부모의 지시가 아닌 아이 스스로 다음 목표를 결정하는 '주도적 몰입' 상태 도달.\n-> 이것이 캘리클럽이 단순한 '놀이터'가 아닌 '성장 엔진'인 본질적 이유입니다.")

    # --- SLIDE 2: THE ENERGY SINK LOGIC ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s2, "에너지 싱크(Energy Sink)와 부모 배터리 보존", "Quantitative Balance: Child's Energy Burn vs. Parental Recovery", 2)
    add_insight_box(s2, 0.5, 2.2, 12.3, 4.2, "Thermodynamic Equilibrium", "1. 에너지 싱크: 5세의 폭발적 활동량을 2시간 내 90% 이상 소진시켜 21시 이전 완전 취침 유도.\n2. 부모 배터리 보존: 동선 500m 내 식사와 휴식이 해결되는 '원스톱 인프라' 활용으로 부모의 물리적 피로 최소화.\n3. 감정적 잉여: 아이의 즐거움이 부모의 죄책감을 덜어주고, 가족 간의 긍정적 대화로 전환되는 감정적 ROI 달성.")

    # --- SLIDE 3: GWANGMYEONG VS THE WORLD ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s3, "광명 클러스터 vs 스타필드: 전략적 우위 분석", "Density & Proximity: Why Gwangmyeong is the Sovereign Choice", 3)
    table = s3.shapes.add_table(4, 3, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.5)).table
    rows = [["Comparison Metric", "Gwangmyeong Cluster (Lotte/IKEA)", "Anseong Starfield (Mega Mall)"], ["Spatial Density", "500m Cluster (High Efficiency)", "3km+ Indoor Trekking (High Fatigue)"], ["Specific Experience", "IT Sports & Discovery Science", "General Commerce & Water Play"], ["Strategic Verdict", "Focus on Deep Engagement", "Focus on Broad Consumption"]]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c); cell.text = val; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(14)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = C_NAVY; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True
    add_insight_box(s3, 0.5, 5.8, 12.3, 1.2, "이번 주말, 동선의 효율성과 아이의 몰입도를 고려한다면 '광명'의 압승입니다.")

    # --- SLIDE 4: THE DECISION TREE ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s4, "Sovereign Decision Tree: 피로도 기반 동선 알고리즘", "Algorithmic Path Selection based on Current Parental State", 4)
    add_insight_box(s4, 0.5, 2.2, 12.3, 4.2, "The Action Algorithm", "1. 부모 배터리 > 80%: [광명동굴 익스플로러] - 163계단의 고난을 딛고 경외감을 선사하십시오.\n2. 부모 배터리 50~80%: [캘리클럽 액티브] - 롯데몰 내에서 모든 에너지를 태우고 쾌적하게 식사하십시오.\n3. 부모 배터리 < 30%: [소올투베이커리 벙커] - 키즈룸 인근 명당에서 커피를 마시며 아이의 놀이를 관조하십시오.\n-> 오늘의 배터리 상태가 오늘의 동선을 결정합니다.")

    # --- SLIDE 5: THE SOVEREIGN VERDICT ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s5, "The Sovereign Verdict: 최종 결론 및 로드맵", "One-Page Strategy for the Perfect Weekend Execution", 5)
    add_insight_box(s5, 0.5, 2.2, 12.3, 4.5, "Final Executive Roadmap", "1. 10:30 AM: 롯데몰 광명 캘리클럽 오픈런 (성취의 심리학 시작).\n2. 12:30 PM: 아브뉴프랑 까몬 '어린이 쌀국수' (영양 및 미식의 균형).\n3. 14:30 PM: 컨디션에 따라 도덕산 출렁다리 또는 광명동굴 선택적 탐험.\n4. Verdict: 이번 주말은 '성취'와 '효율'을 동시에 잡는 [광명역세권 콤팩트 동선]이 정답입니다.\n지금 출발하면 당신은 주말의 주권을 쥐게 될 것입니다.")

    prs.save('Gwangmyeong_Trip_Sovereign_v26.pptx')
    print("V26 Sovereign Strategic (5 Slides) Created. Focused on Essential Logic.")

if __name__ == "__main__":
    create_sovereign_5_ppt_v26()
