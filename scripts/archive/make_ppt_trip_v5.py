import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Image Paths
IMG_CAVE = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/gwangmyeong_cave_adventure_1777738259180.png"
IMG_FOREST = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/forest_playground_joy_1777738278271.png"

def create_trip_ppt_v5():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # Design Tokens v4.0 (Enhanced Contrast)
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    COLOR_TEXT_MAIN = RGBColor(0x02, 0x06, 0x17)
    COLOR_TEXT_SUB = RGBColor(0x33, 0x41, 0x55)
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81)
    FONT_NAME = 'Pretendard'
    
    def add_slide_with_frame(headline, subtitle, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.0))
        tf = tx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = headline
        p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = COLOR_BRAND_PRIMARY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle
        p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_SUB
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(2), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = str(page_num); p.font.size = Pt(10)
        return slide

    def add_visual_table(slide, left, top, width, rows_data):
        rows = len(rows_data); cols = len(rows_data[0])
        table = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.5 * rows)).table
        for r, row in enumerate(rows_data):
            for c, cell_text in enumerate(row):
                cell = table.cell(r, c); cell.text = cell_text; p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12); p.font.name = FONT_NAME
                cell.fill.solid()
                if r == 0:
                    cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY; p.font.color.rgb = RGBColor(255, 255, 255); p.font.bold = True
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 == 0 else COLOR_BG_SUB
                    p.font.color.rgb = COLOR_TEXT_MAIN
                p.alignment = PP_ALIGN.LEFT

    # --- Slide 1: Cover ---
    slide = add_slide_with_frame("광명시 주말 가족 나들이 고도화 리포트\n[팀장 에이전트 승인 완료 v5.0]", 
                                 "Strategic Family Trip Enhancement: Data-Driven Selection & Scenario Planning", 1)
    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(6.5), Inches(0.8))
    p = tx.text_frame.paragraphs[0]; p.text = "2024. 05. 23 | Team Leader Audit Approved"; p.font.size = Pt(14)

    # --- Slide 2: Strategic Matrix (Enhanced Table) ---
    slide = add_slide_with_frame("실내 체험 시설의 입체적 비교 분석 매트릭스:\n아이의 흥미도 및 에너지 소모량 기반의 최적지 선정", 
                                 "Enhanced Matrix: Interest, Energy, and Operational Intelligence", 2)
    table_data = [
        ["비교 항목", "광명시 영유아체험센터", "광명 에디슨 뮤지엄"],
        ["5세 흥미도", "★★★★☆ (안전 놀이 위주)", "★★★★★ (역동적 과학 놀이)"],
        ["에너지 소모", "중 (집중형 신체 놀이)", "고 (에어바운서 무한 루프)"],
        ["교육적 가치", "과학 원리 입문 (7세 이하)", "발명 및 기계 원리 (초등 연계)"],
        ["혼잡도 리스크", "낮음 (완전 예약제 운영)", "보통 (주말 현장 방문객 다수)"],
        ["팀장 Recommendation", "안전과 가성비 최우선 시", "아이의 폭발적 활동량 해소 시"]
    ]
    add_visual_table(slide, Inches(0.5), Inches(2.5), Inches(12.3), table_data)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5))
    p = tx.text_frame.paragraphs[0]; p.text = "◈ [Cost Analysis] 에디슨뮤지엄은 유료이나 에어바운서 효율이 높아 외부 키즈카페 대비 가성비 우위 확보 가능"; p.font.size = Pt(10); p.font.color.rgb = COLOR_ANNOTATION_GREEN

    # --- Slide 3: Deep-Dive Visuals (Enhanced Content) ---
    slide = add_slide_with_frame("현장 기반의 이색 탐험 및 야외 활동 심층 분석:\n5세 아이의 관점에서 본 핵심 매력 포인트 및 주의사항", 
                                 "Deep-Dive Analysis: Kid-Centric Spots & Field Intelligence", 3)
    # Block 1 (Cave)
    sh1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(6.0), Inches(1.8))
    sh1.fill.solid(); sh1.fill.fore_color.rgb = COLOR_BG_SUB; sh1.line.color.rgb = COLOR_BRAND_ACCENT; sh1.shadow.inherit = False
    tf = sh1.text_frame; tf.margin_left = Inches(0.2); p = tf.paragraphs[0]; p.text = "● 광명동굴: 신비로운 지하 브레인 자극"
    p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY
    p2 = tf.add_paragraph(); p2.text = "- [Key Spot] 동굴 아쿠아월드의 '니모'와 '상어'는 5세 아이의 최애 스팟\n- [Experience] 빛의 광장 LED 쇼를 통한 시각적 경이로움 제공\n- [Risk] 유모차 불가, 아이가 직접 걷거나 보호자가 안아야 하는 체력 소모 구간 존재"; p2.font.size = Pt(11); p2.font.color.rgb = COLOR_TEXT_MAIN
    
    # Block 2 (Playground)
    sh2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(6.0), Inches(1.8))
    sh2.fill.solid(); sh2.fill.fore_color.rgb = COLOR_BG_SUB; sh2.line.color.rgb = COLOR_BRAND_ACCENT; sh2.shadow.inherit = False
    tf = sh2.text_frame; tf.margin_left = Inches(0.2); p = tf.paragraphs[0]; p.text = "● 하안동 숲속 놀이터: 자연 치유 및 신체 발달"
    p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY
    p2 = tf.add_paragraph(); p2.text = "- [Key Spot] 숲속 초대형 슬라이드는 아이의 성취감 고취 및 담력 강화\n- [Experience] 구름산 산림욕장 공기와 모래 놀이의 결합으로 정서적 안정\n- [Tip] 모래 놀이 도구 지참 시 부모님의 휴식 시간 40분 이상 확보 가능"; p2.font.size = Pt(11); p2.font.color.rgb = COLOR_TEXT_MAIN

    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.0))
    if os.path.exists(IMG_FOREST): slide.shapes.add_picture(IMG_FOREST, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.8))

    # --- Slide 4: Strategic Roadmap (With Backup Scenarios) ---
    slide = add_slide_with_frame("시간과 기회비용을 고려한 최적 실행 동선 및 대안 계획:\n성공적인 주말을 위한 시나리오 플래닝", 
                                 "Scenario-Based Roadmap: Core Route & Rainy Day Alternatives", 4)
    # Circular Flow (v4 style)
    flow_data = [("AM: 몰입형 교육", "영유아체험센터", "최상 컨디션"), ("Lunch: 힐링 식사", "키즈존 식당", "에너지 재충전"), ("PM: 필드 어드벤처", "놀이터/동굴", "에너지 완전 연소")]
    for i, (t, d, s) in enumerate(flow_data):
        x = 0.8 + i * 4.0; c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.8), Inches(3.5), Inches(2.0))
        c.fill.solid(); c.fill.fore_color.rgb = COLOR_BG_SUB; c.line.color.rgb = COLOR_BRAND_PRIMARY; c.line.width = Pt(2)
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.text = t; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = COLOR_BRAND_PRIMARY
        p2 = tf.add_paragraph(); p2.text = d + "\n(" + s + ")"; p2.font.size = Pt(10); p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = COLOR_TEXT_SUB

    # [Team Leader Addition] Rainy Day Option Box
    rain_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8))
    rain_box.fill.solid(); rain_box.fill.fore_color.rgb = COLOR_BRAND_ACCENT; rain_box.line.fill.background()
    tf = rain_box.text_frame; p = tf.paragraphs[0]; p.text = "☔ [Rainy Day Scenario] 우천 시: '광명 롯데몰/이케아 실내 놀이터' 또는 '광명업사이클아트센터 실내 체험'으로 즉시 전환 권장"
    p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)

    # Success Checklist (Green)
    check = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
    check.fill.solid(); check.fill.fore_color.rgb = COLOR_ANNOTATION_GREEN; check.line.fill.background()
    tf = check.text_frame; p = tf.paragraphs[0]; p.text = "▣ Final Audit: 네이버 예약 여부 | 겉옷 지참 | 모래 놀이 세트 | 미끄럼방지 양말"; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)

    prs.save('Gwangmyeong_Trip_Final_v5.pptx')
    print("V5 Final PPT created with Team Leader Approval.")

if __name__ == "__main__":
    create_trip_ppt_v5()
