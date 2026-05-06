import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

IMG_CALI = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/cali_club_it_sports_vibe_1777742544504.png"
IMG_EDISON = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/edison_museum_discovery_vibe_1777742558123.png"

def create_executive_ppt_v24():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    C_NAVY = RGBColor(0x0F, 0x17, 0x2A); C_BLUE = RGBColor(0x38, 0xBD, 0xF8); C_BG = RGBColor(0xF8, 0xFA, 0xFC)
    
    def apply_style(slide, headline, subtitle, page):
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = C_NAVY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(14); p.font.color.rgb = RGBColor(100, 100, 100)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Executive Strategic Hub v24.0 | Curation Excellence | Page {page}/10"; p.font.size = Pt(10); p.font.color.rgb = RGBColor(180, 180, 180)

    def add_insight_box(slide, x, y, w, h, title, body, color=C_BLUE):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = C_BG; box.line.color.rgb = color; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]; p.text = "🤵 Recommendation: " + title; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = color
        p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(12); p2.font.color.rgb = C_NAVY

    # --- SLIDE 1: COVER ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s1, "광명시 주말 전략 보고서 [v24.0 Executive Summary]", "High-Impact Strategy for Gwangmyeong Weekend Optimization", 1)
    
    # --- SLIDE 2: STRATEGIC LOGIC ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s2, "왜 광명인가? 전략적 가치 분석", "Efficiency, Proximity, and Diverse Experience Density", 2)
    table = s2.shapes.add_table(4, 3, Inches(0.5), Inches(2.2), Inches(12.3), Inches(3.5)).table
    rows = [["Category", "Gwangmyeong Advantage", "Economic Value"], ["Cost (TCO)", "Avg 70k (3P) vs 110k (Tech Park)", "36% Savings"], ["Proximity", "Within 15km from Guro/Gwanak", "Gas/Time Optimization"], ["Diversity", "IT Sports, Science, Nature Cluster", "Multi-Scenario Selection"]]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c); cell.text = val; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(14)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = C_NAVY; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True

    # --- SLIDE 3: MASTER DECISION MATRIX ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s3, "Master Decision Matrix: 오늘의 최적 동선 1초 선택", "Aligning Child Energy vs. Parental Fatigue", 3)
    add_insight_box(s3, 0.5, 2.2, 12.3, 4.0, "The Strategic Compass", "아이의 에너지와 부모의 컨디션을 결합하여 아래 경로 중 하나를 선택하십시오.\n1. Path A (Activity Focus): 5세 아동의 대근육 에너지를 100% 연소하고 싶을 때\n2. Path B (Healing Focus): 부모의 피로도가 극심하여 '프라이빗 휴식'이 필요할 때\n3. Path C (Smart Focus): 단순한 놀이를 넘어 교육적 자극이 필요할 때")

    # --- SLIDE 4: PATH A (ACTIVITY) ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s4, "Path A: 'Extreme Energy Burn' - 신체 활동 극대화", "Top Pick: Cali Club (Lotte Mall) & Kids Bay Park", 4)
    if os.path.exists(IMG_CALI): s4.shapes.add_picture(IMG_CALI, Inches(0.5), Inches(2.2), Inches(6), Inches(4.2))
    add_insight_box(s4, 6.8, 2.2, 6.0, 4.2, "Activity Verdict", "캘리클럽의 RFID 태그액션과 키즈베이파크의 튜브슬라이드 조합은 광명 내 최고의 활동량을 보장합니다.\n- Cali Club Score: 4.8/5.0 (#활동량폭발)\n- Kids Bay Score: 4.7/5.0 (#가족모임최고)")

    # --- SLIDE 5: PATH B (HEALING) ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s5, "Path B: 'Parental Healer' - 정서 회복 및 휴식", "Top Pick: Gureumsan Forest & Salt Bakery (Kids Room)", 5)
    add_insight_box(s5, 0.5, 2.2, 12.3, 4.2, "Healing Verdict", "부모의 휴식이 1순위라면 소올투베이커리의 2층 프라이빗 키즈룸(기차놀이, 칠판 완비)을 요새로 삼으십시오.\n- 팁: 광명동굴 방문 전 이곳에서 '동굴 할인권'을 반드시 챙기세요.\n- 구름산 황톳길: 맨발 걷기로 부모의 혈액순환과 아이의 정서 안정을 동시에 잡습니다.")

    # --- SLIDE 6: PATH C (SMART) ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s6, "Path C: 'Smart Explorer' - 창의력 및 교육 자극", "Top Pick: Edison Museum & Upcycle Art Center", 6)
    if os.path.exists(IMG_EDISON): s6.shapes.add_picture(IMG_EDISON, Inches(6.8), Inches(2.2), Inches(6), Inches(4.2))
    add_insight_box(s6, 0.5, 2.2, 6.0, 4.2, "Explorer Verdict", "에디슨뮤지엄의 11:00 AM 도슨트 해설은 아이에게 발명가라는 새로운 꿈을 심어줍니다.\n- 에어바운서 마감 시각(17:30)을 주의하십시오.\n- 업사이클 아트센터: 환경 보호와 예술적 감각을 동시에 기를 수 있는 광명의 자랑입니다.")

    # --- SLIDE 7: GASTRONOMY MAP ---
    s7 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s7, "Gastronomy Guide: 실패 없는 맛집 3선", "Curation: Child Palate & Parental Ease", 7)
    table = s7.shapes.add_table(4, 3, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.0)).table
    rows = [["Restaurant", "Key Menu / Price", "Decision Point"], ["라라코스트", "Bacon Cream Pasta (10.9k)", "Large Playground / Robot Serving"], ["서울현방", "Donkatsu Set (8.0k)", "Reliable Cleanliness / High-chairs"], ["상상초월갈비", "Galbitang (No-spicy 12.0k)", "Local Flavor / Large Parking Lot"]]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c); cell.text = val; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(13)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = C_NAVY; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True

    # --- SLIDE 8: LOGISTICS CHEAT SHEET ---
    s8 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s8, "Logistics Hub: 주차 및 교통 치트키", "Field Intelligence: Zero Stress Strategy", 8)
    add_insight_box(s8, 0.5, 2.2, 12.3, 4.2, "Tactical Intelligence", "1. 광명동굴: 제2주차장의 '그늘막' 선점이 오후 컨디션을 결정함 (10:30 AM 포화).\n2. 시민체육관: 1시간 무료, 이후 저렴한 요금으로 시내 거점 활용 최적.\n3. 광명역 권역: 주말 14~17시 진입 금지 (이케아/코스트코 인파 주의).")

    # --- SLIDE 9: SAFETY & MEDICAL ---
    s9 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s9, "Safety First: 비상 의료 대응 맵", "Emergency Readiness for Weekend Families", 9)
    add_insight_box(s9, 0.5, 2.2, 12.3, 4.2, "Medical Readiness", "1. 토요일 오전 응급: A 소아과 (12:30 접수 마감 주의).\n2. 일요일 응급: B 연합의원 (18시까지 상시 운영).\n3. 야간 소아과: 중앙대광명병원 응급실 전문의 24H 상주.\n4. 야간 약국: 24시 광명시민약국 활용.")

    # --- SLIDE 10: FINAL VERDICT ---
    s10 = prs.slides.add_slide(prs.slide_layouts[6]); apply_style(s10, "Final Verdict: 당신의 주말을 위한 최종 결론", "Actionable Plan and Expected Family ROI", 10)
    add_insight_box(s10, 0.5, 2.2, 12.3, 4.2, "The Sovereign Verdict", "아이의 성장에 투자하고 싶다면 '에디슨뮤지엄'을, 부모의 휴식이 필요하다면 '소올투베이커리'를 선택하십시오.\n오늘의 추천: [Path A + Gastronomy A] 조합으로 신체 활동과 부모 휴식을 동시에 쟁취하십시오.\n지금 출발하면 광명의 주차장은 당신의 것입니다.")

    prs.save('Gwangmyeong_Trip_Executive_v24.pptx')
    print("V24 Strategic Executive Summary (10 Slides) Created.")

if __name__ == "__main__":
    create_executive_ppt_v24()
