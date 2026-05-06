import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_uhd_ppt():
    prs = Presentation()
    
    # 16:9 Aspect Ratio (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Design Tokens (From v3.0 Final Manual)
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A) # Deep Navy
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)  # Sky Blue
    COLOR_TEXT_MAIN = RGBColor(0x1E, 0x29, 0x3B)    # Slate
    COLOR_TEXT_SUB = RGBColor(0x47, 0x55, 0x69)     # Slate Gray
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)       # Light Grey
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81) # Emerald Green
    
    FONT_NAME = 'Pretendard'
    
    def add_slide_with_frame(headline, subtitle, page_num, dark=False):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background: Pure White (Policy)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Headline (Governance-style Claim)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = headline
        p.font.name = FONT_NAME
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_BRAND_PRIMARY
        p.font.shadow = False
        
        # Subtitle
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = FONT_NAME
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_SUB
        p.font.shadow = False
        
        # Footer
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(2), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        p.text = str(page_num)
        p.font.name = FONT_NAME
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_SUB
        p.font.shadow = False
        
        return slide

    def add_dense_block(slide, left, top, width, title, data, analysis, implication):
        # 2-Pane Style: [Grey Title Box] + [Transparent Content Area]
        
        # 1. Left: Head Title Box (Narrower for balance)
        title_w = Inches(1.5)
        shape_t = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, title_w, Inches(0.6))
        shape_t.fill.solid()
        shape_t.fill.fore_color.rgb = COLOR_BG_SUB
        shape_t.line.color.rgb = RGBColor(200, 200, 200)
        shape_t.shadow.inherit = False
        
        tf_t = shape_t.text_frame
        tf_t.word_wrap = True
        tf_t.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf_t.vertical_anchor = MSO_ANCHOR.TOP # Priority 1: Prevent upward growth
        tf_t.margin_left = Inches(0.1)
        tf_t.margin_right = Inches(0.1)
        
        p = tf_t.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_BRAND_PRIMARY
        p.font.shadow = False
        p.alignment = PP_ALIGN.LEFT
        
        # 2. Right: Body Text Area (Transparent)
        content_w = width - title_w - Inches(0.2)
        shape_c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + title_w + Inches(0.2), top, content_w, Inches(0.6))
        shape_c.fill.background()
        shape_c.line.fill.background()
        
        tf_c = shape_c.text_frame
        tf_c.word_wrap = True
        tf_c.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf_c.vertical_anchor = MSO_ANCHOR.TOP # Priority 1: Prevent upward growth
        
        stages = [
            ("[Data] ", data, COLOR_TEXT_MAIN),
            ("[Analysis] ", analysis, COLOR_TEXT_SUB),
            ("[Strategic Implication] ", implication, COLOR_BRAND_ACCENT)
        ]
        
        for i, (label, text_list, color) in enumerate(stages):
            p = tf_c.paragraphs[i] if i == 0 else tf_c.add_paragraph()
            combined_text = " ".join(text_list) if isinstance(text_list, list) else text_list
            p.text = label + combined_text
            p.font.name = FONT_NAME
            p.font.size = Pt(12) # Strict 12pt policy
            p.font.color.rgb = color
            p.font.shadow = False
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)

    def add_annotation(slide, left, top, text, width=3.0):
        tx = slide.shapes.add_textbox(left, top, Inches(width), Inches(0.5))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "◈ " + text
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_ANNOTATION_GREEN
        p.font.italic = True
        p.font.shadow = False
        p.alignment = PP_ALIGN.LEFT

    # --- Global Parameters ---
    box_w = 6.0
    gap = 0.4
    start_x = (13.333 - (box_w * 2 + gap)) / 2 
    body_top = 2.5 # Priority 1: Avoid subtitle overlap
    vertical_buffer = 0.5 # Safe gap between rows

    # --- Slide 1: Cover (Minimalist White) ---
    slide = add_slide_with_frame("지능형 지식 자산화 및 운영 자동화를 통한\n기업 경쟁력 극대화 전략 제안", 
                                 "AI-Driven Knowledge Framework: Strategic Asset Optimization & Roadmaps", 1)
    txDate = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(10), Inches(0.8))
    tf = txDate.text_frame
    p = tf.paragraphs[0]
    p.text = "2024. 05. 22 | Strategy & Technology Transformation Task Force"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_SUB
    p.font.shadow = False
    p = tf.add_paragraph()
    p.text = "Confidential - Internal Use Only"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.font.shadow = False

    # --- Slide 2: Strategic Necessity ---
    slide = add_slide_with_frame("글로벌 AI 패러다임 시프트에 따른 내부 대응력 격차 진단 및\n데이터 기반 의사결정 체계로의 전략적 전환 필요성",
                                 "Global Trends, Internal Risks, and Economic Value Analysis", 2)
    
    add_dense_block(slide, Inches(start_x), Inches(body_top), Inches(box_w),
                    "글로벌 기술 트렌드 및 시장 동향",
                    ["글로벌 Top-tier 기업의 생성형 AI 도입률 및 R&D 투자 비중 전년 대비 150% 급증(Gartner 기준).", "산업군 내 AI 기반 자동화 프로세스 도입을 통한 운영 효율성(OpEx) 평균 30% 개선 사례 확보."],
                    ["경쟁사 중심의 AI 기술 내재화 가속화에 따른 'AI-Gap' 심화 위험.", "기술 격차를 넘어 운영 프로세스 전반의 생산성 격차로 전이될 수 있는 변곡점 도래."],
                    ["Preemptive Response: 전사적 데이터 자산을 통합 관리할 수 있는 'Scalable AI Platform' 구축 시급."])
    
    add_dense_block(slide, Inches(start_x + box_w + gap), Inches(body_top), Inches(box_w),
                    "내부 운영 효율성 및 데이터 진단",
                    ["정형 데이터의 85%가 사일로(Silo)화되어 있으며 비정형 데이터 활용률은 5% 미만 저조.", "현재 내부 의사결정 프로세스 내 데이터 기반 추론 비중은 12% 수준으로 직관 의존도 높음."],
                    ["Data Silo 현상으로 인한 통합적 통찰(SSOT) 도출 저해 및 의사결정 지연 비용 발생.", "비정형 데이터 미활용은 지식 자산 손실 및 급변하는 환경에서의 Agility 저하 초래."],
                    ["Strategic Pivot: 데이터 사일로를 제거하고 정형/비정형 통합 학습 가능한 구조로 전면 재편 필요."])
    
    # Annotations
    add_annotation(slide, Inches(start_x), Inches(6.5), "PESTEL Analysis: 글로벌 AI 규제 및 거버넌스 대응 체계 반영")
    add_annotation(slide, Inches(start_x + 4.2), Inches(6.5), "Competitor Benchmarking: 업종 내 선도 그룹과의 생산성 Gap")

    # --- Slide 3: Architecture ---
    slide = add_slide_with_frame("데이터 가치 사슬 최적화를 위한 3-Layer 통합 플랫폼 설계 및\n미래 대응형 Future-proof 인프라 구축",
                                 "Target Architecture: Ingestion, Reasoning, and Service Layers", 3)
    
    layers = [
        ("Data Ingestion Layer", 
         ["Multi-source(ERP, CRM, IoT) 기반 통합 Data Lake 및 ETL 파이프라인.", "일일 유입량 500GB 이상의 비정형 데이터 실시간 파싱 및 지능형 벡터화."],
         ["Data Fidelity: 수집부터 임베딩까지의 전 과정 자동화로 데이터 무결성 확보.", "Scalable 아키텍처로 향후 데이터 폭증 시에도 안정적인 시스템 성능 유지."],
         ["Robust Foundation: 어떠한 AI 모델 도입 시에도 즉각 대응 가능한 유연성 확보."]),
        ("Cognitive Reasoning Layer", 
         ["LLM/sLLM Orchestration, RAG 엔진, Vector DB 기반의 정밀 추론 환경.", "Multi-Agent 상호 검증 및 Self-Correction 메커니즘을 통한 신뢰도 극대화."],
         ["Inference Optimization: RAG 기술을 통한 환각 현상 최소화 및 도메인 지식 반영.", "기업 내부 보안 가이드라인을 준수하는 On-premise 하이브리드 추론 모델."],
         ["Intelligence Core: 단순 검색을 넘어 논리적 추론이 가능한 기업용 브레인 역할."])
    ]
    for i, (title, data, analysis, implication) in enumerate(layers):
        add_dense_block(slide, Inches(0.5), Inches(body_top + i*(2.0)), Inches(12.3), title, data, analysis, implication)
    
    # --- Slide 4: Value & Roadmap (Process Layout) ---
    slide = add_slide_with_frame("단계적 AI 생태계 확산을 통한 재무적 가치 창출 및\n지속 가능한 지능형 기업(AI-Native)으로의 체질 개선",
                                 "Financial Impact Analysis and 3-Year Strategic Roadmap", 4)
    
    # Left: Financial Impact
    add_dense_block(slide, Inches(start_x), Inches(body_top), Inches(box_w),
                    "재무적 가치 및 비즈니스 임팩트",
                    ["3년 내 운영 비용(OpEx) 25% 절감 및 업무 처리 속도(Throughput) 2배 향상 목표.", "18개월 BEP 달성 이후 기하급수적 수익성 개선 및 지식 자산 가치 증대 예상."],
                    ["Financial Impact: 비용 절감을 넘어 수요 예측 최적화를 통한 마진 개선 동력 확보.", "데이터 기반 선제적 리스크 관리로 인한 잠재적 손실 비용 회피 효과 극대화."],
                    ["ROI Maximization: 핵심 가치 창출 영역 우선 적용을 통한 투자 대비 성과 조기 가시화."])
    
    # Right: Phased Roadmap (Sequential Blocks)
    roadmap_x = start_x + box_w + gap
    phase_h = 1.35
    phase_gap = 0.15
    phases = [
        ("Phase 1: Foundation", "핵심 데이터 통합 및 PoC 수행. 인프라 구축 및 Quick-win 창출."),
        ("Phase 2: Expansion", "전사 주요 직무별 AI Agent 확산. 업무 자동화 내재화 및 전사 활용."),
        ("Phase 3: Ecosystem", "AI-Native 비즈니스 모델 완성. 파트너 생태계 구축 및 자율 운영 진화.")
    ]
    for i, (p_title, p_desc) in enumerate(phases):
        y_pos = body_top + i * (phase_h + phase_gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(roadmap_x), Inches(y_pos), Inches(box_w), Inches(phase_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_BG_SUB
        shape.line.color.rgb = COLOR_BRAND_ACCENT
        shape.line.width = Pt(1.5)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = "▶ " + p_title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_BRAND_PRIMARY
        p.space_after = Pt(4)
        p2 = tf.add_paragraph()
        p2.text = p_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_TEXT_SUB

    # Final Strategic Summary
    summary_top = 6.7
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(summary_top), Inches(12.3), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_SUB
    shape.line.color.rgb = COLOR_ANNOTATION_GREEN
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = tf.paragraphs[0]
    p.text = " [Strategic Conclusion] 본 프로젝트는 단순 기술 도입을 넘어 'AI-First' 기업으로의 체질 개선을 통한 지속 가능한 경쟁 우위 확보를 지향함"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_ANNOTATION_GREEN
    p.alignment = PP_ALIGN.LEFT

    prs.save('Knowledge_Automation_UHD_Final.pptx')
    print("Final UHD PPT created: Knowledge_Automation_UHD_Final.pptx")

if __name__ == "__main__":
    create_uhd_ppt()
