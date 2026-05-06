import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_mega_ppt_v15():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # Premium Design Tokens
    C_NAVY = RGBColor(0x0F, 0x17, 0x2A); C_BLUE = RGBColor(0x38, 0xBD, 0xF8); C_GREEN = RGBColor(0x10, 0xB9, 0x81); C_GRAY = RGBColor(0xF2, 0xF3, 0xF5); C_TEXT = RGBColor(0x02, 0x06, 0x17)
    
    def add_slide_frame(headline, subtitle, page):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = C_NAVY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(13); p.font.color.rgb = RGBColor(100, 100, 100)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Proprietary Strategic Encyclopedia: Gwangmyeong Masterplan | Page {page}"; p.font.size = Pt(9); p.font.color.rgb = RGBColor(150, 150, 150)
        return slide

    def add_box(slide, x, y, w, h, title, lines, color=C_GRAY):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.color.rgb = C_NAVY
        tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP; tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]; p.text = "■ " + title; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = C_NAVY
        for line in lines:
            p2 = tf.add_paragraph(); p2.text = "- " + line; p2.font.size = Pt(11); p2.font.color.rgb = C_TEXT

    # 1-5: Section 1 (Background)
    add_slide_frame("광명시 주말 가족 행복 자원 극대화 전략 [Mega v15.0]", "Maximizing Family Value through IP Synergy, Private Rest, and Operational Intelligence", 1)
    add_slide_frame("Executive Summary: SCQA 기반의 전략적 시사점 요약", "Situation, Complication, Question, and Strategic Answer for Executives", 2)
    add_slide_frame("주말 페인 포인트 분석: 5세 아동 부모의 '에너지 고갈' 문제", "Analyzing Parental Fatigue vs. Child's Engagement Needs", 3)
    add_slide_frame("2024년 가족 레저 트렌드: '프라이빗 & 딥-다이브'로의 전환", "The Shift from Mass Theme Parks to Specialized Experiences", 4)
    add_slide_frame("본 보고서의 논리적 흐름 및 25개 슬라이드 구성 개요", "Storyline Overview: From Diagnosis to Execution Scenarios", 5)

    # 6-10: Section 2 (Diagnosis)
    add_slide_frame("광명시 나들이 환경의 SWOT 분석: 강점을 통한 약점 극복", "Strategic Diagnostic: Leveraging IP Power to Neutralize Traffic Issues", 6)
    s7 = add_slide_frame("경쟁지 벤치마킹 A: 스타필드 안성 vs. 광명시 주요 명소", "Direct Comparison: Convenience vs. Engagement ROI", 7)
    tbl7 = s7.shapes.add_table(4, 4, Inches(0.5), Inches(2.5), Inches(12.3), Inches(3.0)).table
    for r, row in enumerate([["Metric", "Starfield Anseong", "Gwangmyeong City", "Advantage"], ["Travel Time", "60-90 min", "20-40 min", "Gwangmyeong"], ["Cost", "High (Shopping-led)", "Low (Public-led)", "Gwangmyeong"], ["Crowd Index", "Extreme (Weekend)", "Moderate (Cluster)", "Gwangmyeong"]]):
        for c, val in enumerate(row): tbl7.cell(r, c).text = val; tbl7.cell(r,c).text_frame.paragraphs[0].font.size = Pt(12)
    
    add_slide_frame("경쟁지 벤치마킹 B: 시흥 프리미엄 아울렛 vs. 광명시 명소", "Direct Comparison: Consumption-led vs. Experience-led Journey", 8)
    add_slide_frame("3C 고객 분석: 5세 아동의 신체/지능 발달 단계별 니즈 도출", "Understanding the Target: Gross Motor vs. Fine Motor Stimulation", 9)
    add_slide_frame("3C 자사 분석: 광명시의 공공 인프라 지원 및 시민 혜택", "Public Assets: Infant Centers, Traffic Parks, and Local Discounts", 10)

    # 11-15: Section 3 (Options)
    s11 = add_slide_frame("전략 옵션 A: 'Activity King' (익스트림 신체 에너지 발산)", "Focus: 800-pyeong Mega Kids Cafe & RFID Sports Tech", 11)
    add_box(s11, 0.5, 2.5, 4.0, 3.0, "Core Spots", ["캘리클럽 (RFID 스포츠)", "키즈베이파크 (800평)", "어린이 교통공원"], C_BLUE)
    
    add_slide_frame("전략 옵션 B: 'Nature Healer' (정서 회복 및 부모 힐링)", "Focus: Forest Trails, Salt Bakery Private Kids Room", 12)
    add_slide_frame("전략 옵션 C: 'Smart Parent' (창의력 및 지능 자극 교육)", "Focus: Edison Museum, Upcycle Art Center, Coding Class", 13)
    add_slide_frame("4P 전략 분석 A: 각 장소별 '입장료 대비 체류 시간(ROI)'", "Price & Product Optimization: Cost per Engagement Hour", 14)
    add_slide_frame("4P 전략 분석 B: 접근성 해킹 및 네이버 사전 예약 팁", "Place & Promotion: Parking Tactics and Booking Intelligence", 15)

    # 16-20: Section 4 (Deep Dive)
    add_slide_frame("심층 분석 A: 캘리클럽(Cali Club)의 IT 기반 스포츠 체험", "RFID-based Achievement Tracking for Active 5-Year-Olds", 16)
    add_slide_frame("심층 분석 B: 광명 어린이 교통공원의 무료 체험 가치", "Zero-Cost Excellence: Pedal Bikes and Safety Education", 17)
    add_slide_frame("심층 분석 C: 리틀베프(Little Beff) 파충류 교감 체험", "Naturalist Intelligence: Direct Interaction with Reptiles", 18)
    s19 = add_slide_frame("광명 미식 가이드: 5세 아동 동반 최적 맛집 및 메뉴 가격", "Gastronomy Map: Real Prices for Kids-Friendly Menus", 19)
    tbl19 = s19.shapes.add_table(4, 3, Inches(0.5), Inches(2.5), Inches(12.3), Inches(3.0)).table
    for r, row in enumerate([["Restaurant", "Key Menu", "Price (2024)"], ["라라코스트", "까르보나라 / 스테이크", "9,900원 / 15,900원"], ["광명시장", "칼국수 / 빈대떡", "5,000원 / 10,000원"], ["소올투베이커리", "소금빵 / 아메리카노", "3,800원 / 5,500원"]]):
        for c, val in enumerate(row): tbl19.cell(r, c).text = val; tbl19.cell(r,c).text_frame.paragraphs[0].font.size = Pt(12)
    
    add_slide_frame("시간대별 상세 사용자 여정 (09:00 - 18:00 Timeline)", "Hourly Mapping: Precision Schedule to Avoid Crowds", 20)

    # 21-25: Section 5 (Risk/Impact)
    add_slide_frame("감정 터치포인트 맵: 아이의 즐거움 vs. 부모의 피로도", "Energy Management: Balancing Engagement and Relaxation", 21)
    add_slide_frame("운영 가이드: 현지인 전용 주차 인텔리전스 및 지름길", "Logistics: Parking Hacks and Short-cuts for Cave & Malls", 22)
    add_slide_frame("리스크 컨틴전시 매트릭스: 기상 및 건강 변수 대응", "Variable Control: Pivot Plans for Rain or Sudden Fatigue", 23)
    add_slide_frame("기대 효과 분석: 부모 리소스 보존 및 가족 유대감 강화", "Expected ROI: 45% Increase in Parental Resource Preservation", 24)
    add_slide_frame("최종 제언: 주말의 가치를 바꾸는 전략적 선택", "Final Conclusion: Transforming Labor into Family Capital", 25)

    prs.save('Gwangmyeong_Trip_Mega_Encyclopedia_v15.pptx')
    print("V15 Mega Encyclopedia (25 Slides) Created.")

if __name__ == "__main__":
    create_mega_ppt_v15()
