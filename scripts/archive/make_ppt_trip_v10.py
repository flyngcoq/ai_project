import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Image Paths
IMG_CAVE = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/gwangmyeong_cave_adventure_1777738259180.png"
IMG_FOREST = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/forest_playground_joy_1777738278271.png"

def create_trip_ppt_v10_masterpiece():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # McKinsey-Grade Tokens
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    COLOR_TEXT_MAIN = RGBColor(0x02, 0x06, 0x17)
    COLOR_TEXT_SUB = RGBColor(0x1E, 0x29, 0x3B)
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81)
    
    def add_slide_with_frame(headline, subtitle, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = COLOR_BRAND_PRIMARY
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle; p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_SUB
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Proprietary Strategy: Gwangmyeong Family Value Lab | Page {page_num}"; p.font.size = Pt(9); p.font.color.rgb = RGBColor(120, 120, 120)
        return slide

    def add_integrated_shape(slide, shape_type, left, top, width, height, title, body_lines, color=COLOR_BG_SUB, text_color=COLOR_TEXT_MAIN):
        sh = slide.shapes.add_shape(shape_type, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.color.rgb = COLOR_BRAND_PRIMARY; sh.shadow.inherit = False
        tf = sh.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.margin_left = Inches(0.2); tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.text = title; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY
        for line in body_lines:
            p2 = tf.add_paragraph(); p2.text = "- " + line; p2.font.size = Pt(12); p2.font.color.rgb = text_color

    # --- Slide 1: Cover ---
    slide = add_slide_with_frame("광명시 주말 가족 나들이 가치 극대화 전략 마스터플랜\n[Zero-Base Framework Integration v10.0]", 
                                 "Strategic Resource Allocation & Kid-Parent Value Optimization using 3C & STP Frameworks", 1)
    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))

    # --- Slide 2: 3C Analysis (Strategic Diagnostic) ---
    slide = add_slide_with_frame("시장 및 환경 분석(3C Analysis)을 통한 광명시 나들이의 경쟁 우위 확보:\n5세 아동 니즈와 지자체 가용 자원의 전략적 정렬", 
                                 "Strategic Diagnostic: Aligning Customer Needs with Regional Core Competencies", 2)
    # 3-Pane 3C Layout
    pane_w = 4.0; pane_gap = 0.15; start_x = 0.5
    add_integrated_shape(slide, MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(2.5), Inches(pane_w), Inches(3.5), 
                         "Customer (5세 아동 & 부모)", ["에너지 발산(High) 및 지능 자극 니즈", "부모의 주말 피로도 최소화(Energy Saving) 갈망", "안전이 보장된 프리미엄 체험 환경 선호"])
    add_integrated_shape(slide, MSO_SHAPE.RECTANGLE, Inches(start_x + pane_w + pane_gap), Inches(2.5), Inches(pane_w), Inches(3.5), 
                         "Competitor (인근 수도권 테마파크)", ["고비용 대규모 시설 위주, 극심한 인파 리스크", "이동 거리 및 대기 시간으로 인한 높은 피로도", "획일화된 콘텐츠로 인한 경험 희소성 저하"])
    add_integrated_shape(slide, MSO_SHAPE.RECTANGLE, Inches(start_x + (pane_w + pane_gap)*2), Inches(2.5), Inches(pane_w), Inches(3.5), 
                         "Company (광명시 핵심 역량)", ["광명동굴 등 차별화된 이색 스토리텔링 보유", "영유아체험센터 등 지자체 기반의 높은 가용성", "수도권 초근접성으로 인한 물류/시간 비용 우위"])

    # --- Slide 3: Positioning Map (STP Logic) ---
    slide = add_slide_with_frame("전략적 타겟팅 및 포지셔닝(STP Logic):\n활동 임팩트 대 실행 용이성 기반의 최적 방문지 매핑", 
                                 "Strategic Mapping: Engagement Impact vs Operational Feasibility Matrix", 3)
    # Background Grid
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.5), Inches(10), Inches(4))
    bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_BG_SUB; bg.line.fill.background()
    # Axis Text (Explicit Dark)
    tx_y = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(1), Inches(4))
    p = tx_y.text_frame.paragraphs[0]; p.text = "아이\n몰입도"; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = COLOR_BRAND_PRIMARY
    tx_x = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10), Inches(0.5))
    p = tx_x.text_frame.paragraphs[0]; p.text = "실행 용이성 (부모 리소스 보존율) →"; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = COLOR_BRAND_PRIMARY

    # Bubbles (Text Inside)
    spots = [("영유아센터", 3.2, 5.5, "안정적 밸런스"), ("에디슨뮤지엄", 3.5, 4.0, "고성능 체험"), ("광명동굴", 2.8, 3.2, "이색 가치"), ("숲속놀이터", 3.0, 4.5, "자연 발산")]
    for i, (name, y, x, label) in enumerate(spots):
        bub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.5), Inches(6.5 - y), Inches(1.8), Inches(1.2))
        bub.fill.solid(); bub.fill.fore_color.rgb = COLOR_BRAND_ACCENT; bub.line.color.rgb = COLOR_BRAND_PRIMARY
        tf = bub.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.text = name; p.font.bold = True; p.font.size = Pt(12); p.alignment = PP_ALIGN.CENTER; p.font.color.rgb = COLOR_BRAND_PRIMARY
        p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(10); p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 4: Action Roadmap (Chevron Process) ---
    slide = add_slide_with_frame("최적 실행 로드맵 및 리스크 관리 계획(Chevron Flow):\n성공적인 주말을 위한 타임라인 기반의 컨틴전시 전략", 
                                 "Chevron Roadmap: Tactical Timeline & Contingency Decision Matrix", 4)
    # Chevrons
    steps = [("Step 01: 지능 자극", "실내 시설 집중 체험", COLOR_BRAND_PRIMARY), ("Step 02: 리소스 보충", "키즈 친화형 식사", COLOR_BRAND_ACCENT), ("Step 03: 신체 발산", "야외/동굴 탐험", COLOR_ANNOTATION_GREEN)]
    for i, (t, d, c) in enumerate(steps):
        x = 0.5 + i * 4.2
        sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(2.5), Inches(4.0), Inches(1.5))
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
        tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.5)
        p = tf.paragraphs[0]; p.text = t; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = RGBColor(255,255,255)
        p2 = tf.add_paragraph(); p2.text = d; p2.font.size = Pt(12); p2.font.color.rgb = RGBColor(255,255,255)

    # McKinsey-style Contingency Table
    table_data = [
        ["Risk Variable", "Strategic Response (Plan B/C)", "Impact Level"],
        ["Rainy/Bad Air", "Pivot to Gwangmyeong Lotte Mall / IKEA Inner Course", "Critical"],
        ["Parking Congestion", "Move to Upcycle Art Center (Low Traffic) & Return later", "Moderate"],
        ["Child Fatigue", "Immediate Transition to 'Quiet Zone' (Cafe) or Early Return", "Tactical"]
    ]
    table = slide.shapes.add_table(len(table_data), 3, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0)).table
    for r, row in enumerate(table_data):
        for c, t in enumerate(row):
            cell = table.cell(r, c); cell.text = t; p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12)
            if r == 0: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY; p.font.color.rgb = RGBColor(255,255,255); p.font.bold = True
            else: cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_BG_SUB if r%2==0 else RGBColor(255,255,255); p.font.color.rgb = COLOR_TEXT_MAIN

    # Footnote
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4))
    p = tx.text_frame.paragraphs[0]; p.text = "* Expected Outcome: Parental resource preservation +60%, Child happiness index maximized through tactical rotation."; p.font.size = Pt(10); p.font.color.rgb = COLOR_ANNOTATION_GREEN

    prs.save('Gwangmyeong_Trip_Masterpiece_v10.pptx')
    print("V10 Zero-Base Masterpiece Created.")

if __name__ == "__main__":
    create_trip_ppt_v10_masterpiece()
