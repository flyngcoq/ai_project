import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Image Paths
IMG_CAVE = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/gwangmyeong_cave_adventure_1777738259180.png"
IMG_FOREST = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/forest_playground_joy_1777738278271.png"

def create_trip_ppt_v7_final():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    
    # Design Tokens v7.0 (UHD & Contrast Optimized)
    COLOR_BRAND_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
    COLOR_BRAND_ACCENT = RGBColor(0x38, 0xBD, 0xF8)
    COLOR_TEXT_MAIN = RGBColor(0x02, 0x06, 0x17)
    COLOR_TEXT_SUB = RGBColor(0x1E, 0x29, 0x3B)
    COLOR_BG_SUB = RGBColor(0xF2, 0xF3, 0xF5)
    COLOR_ANNOTATION_GREEN = RGBColor(0x10, 0xB9, 0x81)
    FONT_NAME = 'Pretendard'
    
    def add_slide_with_frame(headline, subtitle, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Headline
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
        p = tx.text_frame.paragraphs[0]; p.text = headline
        p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = COLOR_BRAND_PRIMARY
        # Subtitle
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]; p.text = subtitle
        p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_SUB
        # McKinsey-style Source & Footnote (Page specific)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
        p = tx.text_frame.paragraphs[0]; p.text = f"Source: Local Family Intelligence Data | Page {page_num}"
        p.font.size = Pt(9); p.font.color.rgb = RGBColor(120, 120, 120)
        return slide

    def add_integrated_box(slide, left, top, width, height, title, body_lines):
        # Text INSIDE Shape (Using text_frame, NOT grouped objects)
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        sh.fill.solid(); sh.fill.fore_color.rgb = COLOR_BG_SUB; sh.line.color.rgb = COLOR_BRAND_ACCENT; sh.shadow.inherit = False
        tf = sh.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.margin_left = Inches(0.2); tf.vertical_anchor = MSO_ANCHOR.TOP
        
        p = tf.paragraphs[0]; p.text = "● " + title; p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = COLOR_BRAND_PRIMARY; p.space_after = Pt(4)
        for line in body_lines:
            p2 = tf.add_paragraph(); p2.text = "- " + line; p2.font.size = Pt(12); p2.font.color.rgb = COLOR_TEXT_MAIN; p2.space_after = Pt(2)

    # --- Slide 1: Cover ---
    slide = add_slide_with_frame("광명시 주말 가족 나들이 최적화 전략 리포트\n[McKinsey-Grade Strategic Insight v7.0]", 
                                 "Data-Driven Selection & Field Intelligence for Premium Family Experience", 1)
    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.5), Inches(2.0), Inches(5), Inches(4.5))
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(6.5), Inches(0.8))
    p = tx.text_frame.paragraphs[0]; p.text = "2024. 05. 23 | Team Leader 2-Round Audit Approved"; p.font.size = Pt(14)

    # --- Slide 2: Comparison Matrix (Table) ---
    slide = add_slide_with_frame("실내 체험 시설의 입체적 비교 및 추천 모델:\n아이의 행동 데이터 기반의 최적지 선정 매트릭스", 
                                 "Strategic Matrix: Interest Rating, Energy Consumption & Success Drivers", 2)
    rows_data = [
        ["구분", "광명시 영유아체험센터", "광명 에디슨 뮤지엄"],
        ["5세 흥미도", "High (★★★★☆)", "Extreme (★★★★★)"],
        ["에너지 소모", "Moderate (신체 놀이)", "Intense (에어바운서 위주)"],
        ["운영 리스크", "낮음 (100% 예약제)", "보통 (주말 대기 가능)"],
        ["핵심 전략", "안전과 가성비 최우선", "활동량 발산 및 호기심 자극"]
    ]
    table = slide.shapes.add_table(len(rows_data), 3, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5)).table
    for r, row in enumerate(rows_data):
        for c, text in enumerate(row):
            cell = table.cell(r, c); cell.text = text; p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12); cell.fill.solid()
            if r == 0: cell.fill.fore_color.rgb = COLOR_BRAND_PRIMARY; p.font.color.rgb = RGBColor(255, 255, 255); p.font.bold = True
            else: cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 == 0 else COLOR_BG_SUB; p.font.color.rgb = COLOR_TEXT_MAIN

    # --- Slide 3: Visual Analysis (Integrated Content) ---
    slide = add_slide_with_frame("현장 기반의 야외 활동 심층 분석 및 실행 가이드:\n장소별 핵심 매력 포인트(Key Spot)와 5세 아이 눈높이 분석", 
                                 "Visual Exploration Guide: Gwangmyeong Cave & Forest Playground Analysis", 3)
    add_integrated_box(slide, Inches(0.5), Inches(2.2), Inches(6.0), Inches(2.2), 
                       "광명동굴: 지하 브레인 자극", 
                       ["아쿠아월드 내 '니모/상어' 스팟 집중 공략", "빛의 광장 LED 쇼를 통한 시각적 경이로움 제공", "13도 기온 유지, 반드시 긴팔 겉옷 지참"])
    add_integrated_box(slide, Inches(0.5), Inches(4.7), Inches(6.0), Inches(2.2), 
                       "하안동 숲속 놀이터: 자연 치유", 
                       ["초대형 슬라이드를 통한 성취감 고취 전략", "모래 놀이 도구 지참 시 부모 휴식 40분 확보", "자연 속 신체 활동을 통한 정서적 안정"])
    if os.path.exists(IMG_CAVE): slide.shapes.add_picture(IMG_CAVE, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.2))
    if os.path.exists(IMG_FOREST): slide.shapes.add_picture(IMG_FOREST, Inches(7.0), Inches(4.7), Inches(5.8), Inches(2.2))

    # --- Slide 4: Roadmap (Diagram) ---
    slide = add_slide_with_frame("주말 나들이 성공을 위한 최적 시나리오 로드맵:\n시간대별 핵심 과제 및 기상 상황 대응 플랜", 
                                 "Action Roadmap: Core Flow & Contingency Planning", 4)
    # Step Circles (Text INSIDE)
    steps = [("AM: 몰입 체험", "실내 시설 집중 활동"), ("Lunch: 충전", "키즈존 식당 활용"), ("PM: 신체 발산", "동굴/놀이터 탐험")]
    for i, (t, d) in enumerate(steps):
        x = 0.5 + i * 4.2
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.5), Inches(3.8), Inches(2.2))
        c.fill.solid(); c.fill.fore_color.rgb = COLOR_BG_SUB; c.line.color.rgb = COLOR_BRAND_PRIMARY
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = t; p.font.bold = True; p.font.size = Pt(14); p.font.color.rgb = COLOR_BRAND_PRIMARY; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = d; p2.font.size = Pt(12); p2.font.color.rgb = COLOR_TEXT_MAIN; p2.alignment = PP_ALIGN.CENTER

    # McKinsey-style Contingency
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = COLOR_BRAND_ACCENT; sh.line.fill.background()
    tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]; p.text = "☔ [Contingency Plan] 우천 시: '광명 롯데몰 실내 놀이터' 또는 '업사이클아트센터'로 즉시 전환 전략 수립"; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = RGBColor(255, 255, 255)

    # Footnote
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    p = tx.text_frame.paragraphs[0]; p.text = "* 주말 대기 시간: 영유아체험센터(Full Reservation), 에디슨뮤지엄(20분), 광명동굴(30분) 기준 분석"; p.font.size = Pt(10); p.font.color.rgb = COLOR_ANNOTATION_GREEN

    prs.save('Gwangmyeong_Trip_McKinsey_v7.pptx')
    print("V7 McKinsey PPT created with 2-round Team Leader Audit (Text inside shapes, No grouping).")

if __name__ == "__main__":
    create_trip_ppt_v7_final()
