import sys
from pathlib import Path
sys.path.append(str(Path(__file__).parent.parent))
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    prs = Presentation()
    
    # Define Colors based on Design System
    PURPLE = RGBColor(0x7B, 0x3C, 0xE9)
    NAVY = RGBColor(0x0F, 0x17, 0x2A)
    BLACK = RGBColor(0x1F, 0x1F, 0x1F)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
    CYAN = RGBColor(0x06, 0xB6, 0xD4)

    # Image Paths (using absolute paths provided by the tool)
    DASHBOARD_IMG = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/ai_security_dashboard_futuristic_1778200297333.png"
    LIFECYCLE_IMG = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/ai_lifecycle_governance_concept_1778200311426.png"
    REASONING_IMG = "/Users/flyngcoq/.gemini/antigravity/brain/af9a080a-1e6d-421a-a8ce-ac2d4471de65/ai_agent_reasoning_guardrail_1778200328140.png"

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, NAVY)
        
        # Dashboard Image as background (faded or half-side)
        slide.shapes.add_picture(DASHBOARD_IMG, Inches(5), Inches(0), height=Inches(7.5))
        
        # Title Box with Semi-transparent background (simulated with shape)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5.5), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.5), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4.5), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(18)
        p.font.color.rgb = CYAN
        p.alignment = PP_ALIGN.LEFT

    def add_content_slide(title_text, content_items, image_path=None, table_data=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, WHITE)
        
        # Header Line & Title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.1), Inches(0.5))
        line.fill.solid()
        line.fill.fore_color.rgb = PURPLE
        line.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(8), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = NAVY
        
        # Main Layout: If image, split 50/50. If table, handle separately.
        content_width = Inches(9) if not image_path else Inches(4.5)
        top = 1.2
        
        # Image Handling
        if image_path:
            slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.5), width=Inches(4.3))
            
        # Text Content
        for item in content_items:
            if isinstance(item, str):
                if item.startswith("###"): # Subheader style
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), content_width, Inches(0.4))
                    tf = box.text_frame
                    p = tf.paragraphs[0]
                    p.text = item.replace("###", "").strip()
                    p.font.bold = True
                    p.font.size = Pt(16)
                    p.font.color.rgb = PURPLE
                    top += 0.4
                else:
                    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), content_width, Inches(0.5))
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "• " + item
                    p.font.size = Pt(13)
                    p.font.color.rgb = BLACK
                    top += 0.35
            elif isinstance(item, list):
                for sub in item:
                    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), content_width - Inches(0.3), Inches(0.3))
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "  - " + sub
                    p.font.size = Pt(11)
                    p.font.color.rgb = GRAY
                    top += 0.28
                top += 0.1
        
        # Table Handling
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0])
            left = Inches(0.5)
            top_table = top + 0.2
            width = Inches(9)
            height = Inches(0.3 * rows)
            
            table = slide.shapes.add_table(rows, cols, left, top_table, width, height).table
            for r in range(rows):
                for c in range(cols):
                    cell = table.cell(r, c)
                    cell.text = str(table_data[r][c])
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    if r == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = PURPLE
                        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
                        cell.text_frame.paragraphs[0].font.bold = True

    # --- PPT Slides Generation ---

    # 1. Title
    add_title_slide("Next-Gen AI Security Strategy", "Beyond Gateway: 에이전틱 시대를 위한 초격차 보안 거버넌스")

    # PART 1. 전략적 배경
    add_content_slide("PART 1. 개요 및 전략적 배경", [
        "### 전략적 개요",
        "단순한 'AI 입구 보안(Gateway)'을 넘어선 'Full-Stack AI Governance' 구축",
        "AI 모델의 생성-운영-폐기(Lifecycle) 전 과정에서의 데이터 안전성 및 추론 무결성 보장",
        "2026년 에이전틱 AI(Agentic AI) 시대의 핵심 보안 인프라",
        "### 전략 목표 (Strategic Objectives)",
        "Visibility: 에이전트 사고 과정(CoT) 및 데이터 흐름의 완벽 가시성",
        "Control: 자율형 AI 일탈 행위에 대한 실시간 'Kill-Switch'",
        "Trust: 증명 가능한 보안 이력을 통한 AI 서비스 수용성 제고",
        "Efficiency: 규제 대응 자동화를 통한 보안 조직 운영 효율 극대화"
    ], image_path=DASHBOARD_IMG)

    add_content_slide("글로벌 표준 부합성: Gartner AI TRiSM", [
        "### AI TRiSM (Trust, Risk, and Security Management)",
        "Gartner가 제시한 AI 거버넌스의 표준 프레임워크를 100% 구현",
        "1. Trust (신뢰성): 설명 가능하고 투명한 AI (CoT 감시, Compliance XAI)",
        "2. Risk (리스크): 데이터 유출 및 모델 환각 선제적 방어 (Semantic DLP, DSPM)",
        "3. Security (보안): 런타임 공격 및 데이터 오염 실시간 차단 (Agent-SPM, Anti-Poisoning)",
        "### Strategic Alignment",
        "단순 정책 수립을 넘어선 '실시간 기술적 강제(Runtime Enforcement)' 기반의 초격차 거버넌스 제공"
    ])

    add_content_slide("시장의 페인포인트 및 위협 환경", [
        "입구 보안(North-South)의 한계: 내부 에이전트 간 통신(East-West) 보안 부재",
        "의미론적 우회(Semantic Evasion): 문맥을 이용한 지능형 프롬프트 인젝션 급증",
        "AI 블랙박스 리스크: 의사결정 근거 불투명으로 인한 책임 소재 불분명",
        "섀도우 AI & 자산 방치: 관리되지 않는 모델 및 퇴역 모델 내 데이터 유출(Memory Leakage)",
        "규제 피로도: 금감원/KISA/EU AI Act 등 복잡해지는 글로벌 규제 대응 한계"
    ])

    edge_table = [
        ["구분", "AI Gateway (Red Ocean)", "AI-DSPM (Growing)", "Agent-SPM (Blue Ocean)"],
        ["관점", "통로(Perimeter) 중심", "데이터(Storage) 중심", "행위(Behavior) 중심"],
        ["핵심 가치", "입구 차단, 속도 제한", "정보 유출 방지, 분류", "AI의 자율성 통제 및 책임"],
        ["진입 장벽", "낮음 (오픈소스 다수)", "중간 (기술 가시성)", "높음 (LLM 추론 분석)"],
        ["금융권 반응", "기본 인프라로 인식", "컴플라이언스 대응용", "AX 전면 도입을 위한 필수재"]
    ]
    add_content_slide("경쟁 우위 분석 (2026 Competitive Edge)", [
        "기존 Gateway 시장은 이미 레드오션화 되었으며, '행위(Behavior)' 중심의 Agent-SPM이 차세대 블루오션으로 부상"
    ], table_data=edge_table)

    # PART 2. 서비스 아키텍처
    add_content_slide("PART 2. Safe AI 3중 방어 체계 아키텍처", [
        "### 4단계 통합 보안 레이어",
        "LAYER 1: 인바운드 게이트웨이 (DLP, 쿼터, 서킷브레이커)",
        "LAYER 2: 에이전트 추론 감시 (CoT 인터셉트, Shadow Reasoning)",
        "LAYER 3: 데이터 및 RAG 보안 (PII 마스킹, 안티 포이즈닝)",
        "LAYER 4: 거버넌스 및 자산 생애주기 (자동 보고, 완전 파기)",
        "### 기술적 차별성",
        "East-West 통신 보안: 에이전트 간(A2A) 및 외부 도구 호출 실시간 모니터링",
        "증명 가능한 지능: 추론 근거 전수 기록 및 규제 리포팅 자동화"
    ], image_path=REASONING_IMG)

    # PART 3. 상세 정책
    add_content_slide("핵심 기술: CoT 실시간 감시 및 Shadow Reasoning", [
        "### 왜 입구(Prompt) 보안만으로는 부족한가?",
        "간접 인젝션(Indirect Injection): 정상 요청 속에 숨겨진 악의적 지시는 추론 과정에서만 포착 가능",
        "논리적 비약 방지: AI가 스스로 보안 절차를 생략하거나 규정을 우회하는 결정 차단",
        "### 3단계 가드레일 메커니즘",
        "1. 강제적 사고 노출 (Enforced Disclosure): <thought> 태그 내 추론 텍스트화",
        "2. 실시간 토큰 인터셉트: 생성 즉시 게이트웨이 레벨 가로채기",
        "3. 그림자 추론 검증 (Shadow Reasoning): 별도 경량 sLLM을 통한 실시간 스캔 및 Kill-Switch"
    ])

    add_content_slide("한국형 특화 보안 및 성능 전략", [
        "### K-Compliance Strategy",
        "망분리 5호 예외 지원: U+ 인프라 내 PaaS 래핑 구조를 통한 샌드박스 우회",
        "행정 자동화: 금감원 가이드라인 맞춤형 '보안 사고 방어 일지' 원클릭 생성",
        "### Performance Optimization",
        "스트리밍 비동기 분석: TTFT를 최소화하는 실시간 병렬 분석",
        "초경량 보안 가디언: 2B 이하 sLLM 배치로 연산 부하 5% 미만 유지"
    ])

    add_content_slide("AI 자산 생애주기 및 파기 (AI-Sanitizer)", [
        "### 모델 퇴역 및 자산 소멸 전략",
        "문제: 모델 내 학습 데이터의 '메모리 현상(Memory Leakage)' 및 역공학 위협",
        "### 핵심 솔루션: AI-Sanitizer",
        "멀티레이어 와이핑(Wiping): GPU 메모리, 벡터 DB, 가중치 파일 완전 삭제",
        "보안 파기 증명서: KISA 가이드라인 준수 증빙 리포트 자동 생성",
        "Lifecycle Dashboard: 전사 AI 자산 인벤토리 및 보안 드리프트 상시 감시"
    ], image_path=LIFECYCLE_IMG)

    add_content_slide("상품화 로드맵 (5대 라인업)", [
        "1. AI-Sanitizer: 자산 완전 파기 및 증명 솔루션",
        "2. AI Lifecycle Governance: 전사 자산 관리 대시보드",
        "3. Model Memory Leakage Scanner: 유출 리스크 자동 감사기",
        "4. AI Compliance-as-a-Service: 글로벌 규제 대응 자동화 패키지",
        "5. AI Cyber Insurance Linkage: 보안 점수 기반 보험료 할인 연계 상품"
    ])

    # PART 4. 세일즈 로직
    add_content_slide("PART 4. CISO 설득 논리 및 성공 시나리오", [
        "### Key Messages",
        "Accountability: 사고 책임 소재 명확화 및 법적 면책 근거 제공",
        "Enabler: 혁신의 방해자가 아닌 '안전한 AX의 조력자'로의 포지셔닝",
        "Asset Protection: 직접적 재무 손실(환각/오작동) 차단",
        "### 신한은행 '평화로운 하루' 시나리오",
        "오전: AI-SPM이 API Key 노출 탐지 및 즉시 격리",
        "오후: Agent-SPM이 에이전트의 투자 성향 조작 의도 포착 및 차단",
        "마감: 금감원 제출용 50페이지 실사 리포트 자동 생성 및 정시 퇴근"
    ])

    # PART 5. 로드맵
    roadmap_table = [
        ["단계", "핵심 타겟", "수요(Urgency)", "수익성(Profit)"],
        ["Phase 1 (단기)", "망분리 대응 금융권", "최상 (필수)", "중간"],
        ["Phase 2 (중기)", "AX 전면 도입 대기업", "높음 (A2A 리스크)", "상 (기술 프리미엄)"],
        ["Phase 3 (장기)", "생태계 및 리스크 관리", "보통", "최상 (보험 연계)"]
    ]
    add_content_slide("서비스 진화 로드맵 요약", [
        "단기적으로 규제 대응(망분리) '미끼 상품'으로 진입하여, 중기적 초격차 기술(SPM)로 해자를 구축하고, 장기적 금융 생태계(보험/파기)로 수익 극대화"
    ], table_data=roadmap_table)

    save_path = "/Users/flyngcoq/AI_Project/40_Projects/Work/Safe_AI/Next_Gen_AI_Security_Strategy.pptx"
    prs.save(save_path)
    print(f"✅ PPT 생성 완료: {save_path}")

if __name__ == "__main__":
    create_ppt()
