import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_pro_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    C_BG = RGBColor(10, 15, 36) # Deep Navy/Black
    C_TEXT = RGBColor(240, 245, 250) # Near White
    C_ACCENT_CYAN = RGBColor(0, 255, 204) # Neon Cyan
    C_ACCENT_PINK = RGBColor(255, 51, 153) # Neon Pink
    C_MUTED = RGBColor(100, 110, 130)

    def apply_base(slide, title, subtitle):
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()
        
        # Header line
        hline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.05))
        hline.fill.solid()
        hline.fill.fore_color.rgb = C_ACCENT_CYAN
        hline.line.fill.background()
        
        # Title
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.333), Inches(0.8))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title.upper()
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = C_TEXT
        
        # Subtitle
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.4))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = C_MUTED
        
        # Footer
        tx3 = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.3))
        p3 = tx3.text_frame.paragraphs[0]
        p3.text = "PRO EDITION // GWANGMYEONG STRATEGIC ASSET // V29.0"
        p3.font.size = Pt(9)
        p3.font.color.rgb = C_MUTED
        
    def add_card(slide, x, y, w, h, title, body, is_risk=False):
        accent = C_ACCENT_PINK if is_risk else C_ACCENT_CYAN
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(20, 25, 45) # Slightly lighter than bg
        card.line.color.rgb = accent
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = accent
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = C_TEXT
        p2.space_before = Pt(10)

    # SLIDE 1: COVER
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base(s1, "THE SOVEREIGN THESIS", "Redefining the 5-Year-Old Weekend Experience in Gwangmyeong")
    add_card(s1, 0.5, 2.5, 12.333, 4.0, "EXECUTIVE SUMMARY", "대부분의 부모는 주말을 '소비'하지만, 우리는 주말을 '지배'합니다.\n이 보고서는 준서(5세)의 폭발적인 에너지를 '성취의 인지'로 변환하고,\n부모의 체력 고갈을 '전략적 휴식'으로 방어하는 궁극의 메커니즘입니다.\n\n광명 클러스터의 고밀도 인프라를 활용하여,\n오직 승리감만이 남는 'High-ROI' 주말을 선사합니다.")

    # SLIDE 2: EMOTIVE ROADMAP
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base(s2, "THE EMOTIVE ROADMAP", "Balancing Child Achievement with Parental Burnout")
    add_card(s2, 0.5, 2.5, 6.0, 4.0, "JUN-SEO'S VICTORY ENGINE", "아이의 목표: '나도 할 수 있어!'라는 주도적 효능감 획득.\n\n- 캘리클럽의 RFID 태그액션은 단순 신체놀이가 아닌 도파민 최적화 과정입니다.\n- 키즈베이파크(무한 발산) 대신 캘리클럽을 선택함으로써,\n'놀이터'를 '성취의 훈련장'으로 격상시킵니다.")
    add_card(s2, 6.833, 2.5, 6.0, 4.0, "PARENTAL BATTERY DEPLETION", "부모의 과제: 물리적 피로의 통제와 휴식의 확보.\n\n- 동굴 탐험의 163계단은 압도적 경외감을 주지만 뼈아픈 체력 소모를 요구합니다.\n- 동선이 끊기지 않는 '원스톱(Lotte-IKEA-Avenue France)' 전략만이\n월요일 출근의 재앙을 막을 유일한 방패입니다.", is_risk=True)

    # SLIDE 3: ACTIVITY SINGULARITY
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base(s3, "THE ACTIVITY SINGULARITY", "Tactical Execution: Cali Club & Gastronomy Fusion")
    add_card(s3, 0.5, 2.5, 12.333, 4.0, "10:30 AM OPEN-RUN & CAMMON SYNERGY", "1. 오픈런의 절대성: 주말 11시 이후 캘리클럽 대기 40팀 돌파. 10:30분 태블릿 등록은 협상의 대상이 아님.\n\n2. 까몬(Cammon) 연계: 극도의 에너지 소비 후, 아브뉴프랑 까몬의 '5,000원 어린이 쌀국수'로 단백질을 즉각 보충.\n\n3. 이케아 백업: 동생(12개월 미만) 동반 시, 이케아 레스토랑의 무료 병 이유식 제도를 적극 활용하여 변수를 차단.")

    # SLIDE 4: CAVE & SANCTUARY
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base(s4, "THE CAVE & SANCTUARY PARADOX", "Risk Management: 163 Stairs vs. 50% Discount Recovery")
    add_card(s4, 0.5, 2.5, 6.0, 4.0, "GWANGMYEONG CAVE CRISIS", "리스크: 유모차 반입 전면 금지. 12도의 한기. 지하로 꽂히는 163계단.\n\n대응: 아기띠 지참 필수, 얇은 바람막이 사전 장착.\n부모의 체력이 80% 이상일 때만 강행할 것.", is_risk=True)
    add_card(s4, 6.833, 2.5, 6.0, 4.0, "SALT BAKERY RECOVERY", "전술: 광명동굴 방문 '전'에 소올투베이커리 선방문 필수.\n-> 빵 구매 영수증으로 동굴 입장권 50% 즉시 할인 획득.\n\n명당 사수: 2층 키즈룸 바로 옆 창가 자리.\n부모는 아이를 방목하며 완벽한 '침묵의 충전'을 수행함.")

    # SLIDE 5: LOGISTICAL VERDICT
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_base(s5, "THE MASTER LOGISTICAL VERDICT", "Zero-Stress Mobility & Emergency Infrastructure")
    add_card(s5, 0.5, 2.5, 12.333, 4.0, "EXECUTIVE DIRECTIVES", "1. 11:15 AM Singularity: 이 시각 제2주차장은 매진됨. KTX A~D 사전 예약 혹은 석수3동 둔치 무료주차장(Plan B) 숙지.\n\n2. Medical Shield: 준소아과(철산) 달빛어린이병원 대기, 중앙대광명병원 24H 전문의 위치 파악 완료.\n\n최종 결론: 캘리클럽에서 승리하고, 소올투에서 충전하라.\n오직 치밀하게 계획된 자만이 주말의 완벽한 주권을 차지할 것이다.")

    prs.save('Gwangmyeong_Trip_Pro_Edition_v29.pptx')
    print("V29 Pro Edition (No-Loop, Dark Theme, 5 Slides) Created Successfully.")

if __name__ == "__main__":
    create_pro_ppt()
