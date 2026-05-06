from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_slide_2(file_path):
    try:
        prs = Presentation(file_path)
        if len(prs.slides) < 2:
            print("Error: The PPT has less than 2 slides.")
            return

        slide = prs.slides[1] # 2nd slide
        print(f"--- Analysis of Slide 2 in {file_path} ---")
        
        for i, shape in enumerate(slide.shapes):
            print(f"\n[Shape {i+1}]")
            print(f"  Name: {shape.name}")
            print(f"  Type: {shape.shape_type}")
            print(f"  Pos: Left={shape.left.inches:.2f}\", Top={shape.top.inches:.2f}\"")
            print(f"  Size: Width={shape.width.inches:.2f}\", Height={shape.height.inches:.2f}\"")
            
            if shape.has_text_frame:
                texts = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
                if texts:
                    print(f"  Text: {' / '.join(texts)}")
            
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.fill.type is not None:
                    try:
                        color = shape.fill.fore_color.rgb
                        print(f"  Color: RGB({color[0]}, {color[1]}, {color[2]})")
                    except:
                        pass

    except Exception as e:
        print(f"Error analyzing PPT: {e}")

if __name__ == "__main__":
    analyze_slide_2('Project_Knowledge_Automation_Fancy.pptx')
